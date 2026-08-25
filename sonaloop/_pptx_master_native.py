"""Adapt generated slides to the customer's native PowerPoint placeholders and grid."""
from __future__ import annotations

from collections import Counter
from hashlib import sha256
import io
from pathlib import Path
import re
from typing import Any
import zipfile


_TITLE_TYPES = {"title", "center_title"}
_BODY_TYPES = {"body", "object", "subtitle"}
_STRUCTURAL_KINDS = {"cover", "title", "agenda", "section", "canvas-section", "closing"}


def _placeholder_type(shape) -> str:
    raw = shape.placeholder_format.type
    return str(getattr(raw, "name", raw)).casefold()


def _placeholders(slide, types: set[str]) -> list[Any]:
    return [shape for shape in slide.placeholders if _placeholder_type(shape) in types]


def _largest(shapes: list[Any]):
    return max(shapes, key=lambda shape: int(shape.width) * int(shape.height)) if shapes else None


def _remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _generated_shapes(slide) -> list[Any]:
    return [shape for shape in slide.shapes if not shape.is_placeholder]


def _clear_generated(slide) -> None:
    for shape in list(_generated_shapes(slide)):
        _remove_shape(shape)


def _set_placeholder(shape, paragraphs: list[str], *, secondary_size: float | None = None) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt

    clean = [str(text).strip() for text in paragraphs if str(text or "").strip()]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    # Keep the customer's placeholder geometry and theme font, but let PowerPoint shrink
    # unusually long localized copy instead of spilling outside the box. This writes native
    # <a:normAutofit>; it does not bake in a font or resize the master placeholder.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if not clean:
        return
    for index, text in enumerate(clean):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = text
        # Layouts frequently carry inherited bullets/automatic numbering.  Sonaloop's
        # strings already contain their visible numbering, so explicitly disable both.
        ppr = paragraph._p.get_or_add_pPr()
        bullet_tags = {qn("a:buAutoNum"), qn("a:buChar"), qn("a:buBlip"), qn("a:buNone")}
        for child in list(ppr):
            if child.tag in bullet_tags:
                ppr.remove(child)
        ppr.insert(0, OxmlElement("a:buNone"))
        if index and secondary_size:
            for run in paragraph.runs:
                run.font.size = Pt(secondary_size)


def _populate_cover(slide, spec: dict[str, Any]) -> bool:
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    if title is None:
        return False
    _clear_generated(slide)
    _set_placeholder(title, [spec.get("title") or ""])
    subtitles = _placeholders(slide, {"subtitle"})
    if subtitles:
        _set_placeholder(_largest(subtitles), [spec.get("subtitle") or "", spec.get("lead") or ""],
                         secondary_size=14)
    body = _placeholders(slide, {"body", "object"})
    if body:
        # Native covers often reserve a narrow footer beside fixed brand artwork. Audience,
        # duration and date belong in notes unless the plan explicitly supplies compact copy;
        # concatenating all three is what caused real customer masters to overlap their logo.
        _set_placeholder(_largest(body), [str(
            spec.get("native_meta") or spec.get("eyebrow") or ""
        ).upper()])
    pictures = _placeholders(slide, {"picture"})
    picture = _largest(pictures)
    image = spec.get("image")
    picture_filled = False
    if picture is not None and isinstance(image, (str, Path)) and Path(image).exists():
        try:
            picture.insert_picture(str(image))
            picture_filled = True
        except Exception:
            pass
    # Empty placeholders are invisible in slide show/PDF but PowerPoint displays large editing
    # prompts. Keep only the filled picture and text placeholders on the generated cover.
    for candidate in list(pictures):
        if candidate is not picture or not picture_filled:
            _remove_shape(candidate)
    for candidate in list(_placeholders(slide, {"subtitle", "body", "object"})):
        if not candidate.has_text_frame or not candidate.text_frame.text.strip():
            _remove_shape(candidate)
    return True


def _populate_agenda(slide, spec: dict[str, Any]) -> bool:
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    body = _largest(_placeholders(slide, {"body", "object"}))
    if title is None or body is None:
        return False
    _clear_generated(slide)
    _set_placeholder(title, [spec.get("heading") or "Contents"])
    _set_placeholder(body, [f"{index:02d}   {item}"
                            for index, item in enumerate(spec.get("items") or [], 1)])
    return True


def _populate_section(slide, spec: dict[str, Any]) -> bool:
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    if title is None:
        return False
    _clear_generated(slide)
    _set_placeholder(title, [spec.get("title") or "", spec.get("subtitle") or ""],
                     secondary_size=14)
    body = _largest(_placeholders(slide, {"body", "object"}))
    if body is not None:
        _set_placeholder(body, [spec.get("num") or ""])
    return True


def _populate_closing(slide, spec: dict[str, Any]) -> bool:
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    if title is None:
        return False
    _clear_generated(slide)
    secondary = [spec.get("text") or "", spec.get("contact") or ""]
    _set_placeholder(title, [spec.get("title") or "", *secondary], secondary_size=14)
    return True


def _populate_image(slide, spec: dict[str, Any]) -> bool:
    image = spec.get("image")
    picture = _largest(_placeholders(slide, {"picture"}))
    if picture is None or not isinstance(image, (str, Path)) or not Path(image).exists():
        return False
    _clear_generated(slide)
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    if title is not None:
        _set_placeholder(title, [str(spec.get("heading") or "")])
    try:
        picture.insert_picture(str(image))
    except Exception:
        return False
    body = _largest(_placeholders(slide, {"body", "object", "subtitle"}))
    if body is not None:
        _set_placeholder(body, [spec.get("caption") or ""])
    return True


def _scale_text(shape, factor: float) -> None:
    if factor >= 0.98:
        return
    frames = []
    if getattr(shape, "has_text_frame", False):
        frames.append(shape.text_frame)
    if getattr(shape, "has_table", False):
        frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
    for frame in frames:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    run.font.size = int(run.font.size * factor)


def _content_target(slide):
    bodies = _placeholders(slide, {"body", "object"})
    if not bodies:
        return None
    left = min(int(shape.left) for shape in bodies)
    top = min(int(shape.top) for shape in bodies)
    right = max(int(shape.left + shape.width) for shape in bodies)
    bottom = max(int(shape.top + shape.height) for shape in bodies)
    return left, top, right - left, bottom - top


def _adapt_content(slide, spec: dict[str, Any], slide_width: int, slide_height: int) -> bool:
    from pptx.util import Inches

    heading = str(spec.get("heading") or "").strip()
    title = _largest(_placeholders(slide, _TITLE_TYPES))
    generated = list(_generated_shapes(slide))
    if heading and title is not None:
        # A native master controls title geometry.  Prefix numbers are a Sonaloop-default
        # decoration and can sit beneath customer artwork or outside a narrow title placeholder.
        _set_placeholder(title, [heading])
        header_bottom = Inches(1.52)
        for shape in list(generated):
            if int(shape.top + shape.height) <= header_bottom:
                _remove_shape(shape)
        generated = list(_generated_shapes(slide))

    target = _content_target(slide)
    if target is None or not generated:
        return bool(title and heading)
    source_left = Inches(0.7)
    source_top = Inches(1.55 if heading else 0.5)
    source_width = slide_width - Inches(1.4)
    source_height = slide_height - source_top - Inches(0.5)
    target_left, target_top, target_width, target_height = target
    if min(target_width, target_height, source_width, source_height) <= 0:
        return bool(title and heading)
    sx = target_width / source_width
    sy = target_height / source_height
    for shape in generated:
        rel_left = int(shape.left) - source_left
        rel_top = int(shape.top) - source_top
        shape.left = int(target_left + rel_left * sx)
        shape.top = int(target_top + rel_top * sy)
        shape.width = max(1, int(shape.width * sx))
        shape.height = max(1, int(shape.height * sy))
        if min(sx, sy) < 0.78:
            _scale_text(shape, max(0.62, min(sx, sy)))
    # Generated visual slides use the placeholder only as a layout grid. Once their editable
    # native shapes are fitted, remove unused body/picture prompts from the editable canvas.
    for placeholder in list(_placeholders(slide, _BODY_TYPES | {"picture"})):
        _remove_shape(placeholder)
    return True


def apply_master_native_layout(slide, spec: dict[str, Any], *,
                               slide_width: int, slide_height: int) -> dict[str, Any]:
    """Make one generated slide use its selected layout's native placeholders and grid."""
    kind = str(spec.get("kind") or "content").casefold()
    if kind in {"cover", "title"}:
        native = _populate_cover(slide, spec)
    elif kind == "agenda":
        native = _populate_agenda(slide, spec)
    elif kind in {"section", "canvas-section"}:
        native = _populate_section(slide, spec)
    elif kind == "closing":
        native = _populate_closing(slide, spec)
    elif kind == "image" and _populate_image(slide, spec):
        native = True
    else:
        native = _adapt_content(slide, spec, slide_width, slide_height)
    return {
        "kind": kind,
        "layout": str(getattr(slide.slide_layout, "name", "") or ""),
        "native": bool(native),
        "filled_placeholders": sum(
            1 for shape in slide.placeholders
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()),
    }


def _package_hashes(data: bytes) -> dict[str, str]:
    prefixes = ("ppt/slideMasters/", "ppt/slideLayouts/", "ppt/theme/", "ppt/media/")
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        return {name: sha256(archive.read(name)).hexdigest()
                for name in archive.namelist() if name.startswith(prefixes)}


def enforce_text_color_rules(data: bytes, *, forbidden: list[str], replacement: str
                             ) -> tuple[bytes, int]:
    """Replace explicit forbidden *text* colors without touching shapes/master parts.

    A workspace style reference may establish a semantic rule such as "yellow is
    a marker, never typography" even when that yellow technically passes a
    contrast threshold on another canvas.  The OOXML walk is intentionally
    limited to text run properties inside generated slide instances.
    """
    from xml.etree import ElementTree as ET

    banned = {str(value or "").lstrip("#").upper() for value in forbidden
              if re.fullmatch(r"#?[0-9A-Fa-f]{6}", str(value or ""))}
    replacement = str(replacement or "").lstrip("#").upper()
    if not banned or not re.fullmatch(r"[0-9A-Fa-f]{6}", replacement):
        return data, 0
    namespace = "http://schemas.openxmlformats.org/drawingml/2006/main"
    changed = 0
    source = io.BytesIO(data)
    output = io.BytesIO()
    with zipfile.ZipFile(source, "r") as incoming, zipfile.ZipFile(output, "w") as outgoing:
        for item in incoming.infolist():
            payload = incoming.read(item.filename)
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", item.filename):
                root = ET.fromstring(payload)
                for props_tag in ("rPr", "defRPr", "endParaRPr"):
                    for props in root.findall(f".//{{{namespace}}}{props_tag}"):
                        for color in props.findall(
                                f"./{{{namespace}}}solidFill/{{{namespace}}}srgbClr"):
                            if str(color.attrib.get("val") or "").upper() in banned:
                                color.set("val", replacement)
                                changed += 1
                if changed:
                    payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            outgoing.writestr(item, payload)
    return output.getvalue(), changed


def inspect_rendered_master_deck(data: bytes, master_template: bytes) -> dict[str, Any]:
    """Return a content-free QA record proving master retention and safe slide geometry."""
    from pptx import Presentation
    from ._pptx_master import master_color_map, master_palette

    prs = Presentation(io.BytesIO(data))
    source_hashes = _package_hashes(master_template)
    output_hashes = _package_hashes(data)
    exact = sum(output_hashes.get(name) == digest for name, digest in source_hashes.items())
    out_of_bounds = []
    fonts: Counter[str] = Counter()
    filled = 0
    for index, slide in enumerate(prs.slides, 1):
        filled += sum(1 for shape in slide.placeholders
                      if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
        for shape in slide.shapes:
            if not shape.is_placeholder and (
                    int(shape.left) < 0 or int(shape.top) < 0
                    or int(shape.left + shape.width) > int(prs.slide_width)
                    or int(shape.top + shape.height) > int(prs.slide_height)):
                out_of_bounds.append({"slide": index, "shape": str(shape.name)})
            frames = [shape.text_frame] if getattr(shape, "has_text_frame", False) else []
            if getattr(shape, "has_table", False):
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
            for frame in frames:
                for paragraph in frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts[str(run.font.name)] += 1
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        xml = b"".join(archive.read(name) for name in archive.namelist()
                       if re.fullmatch(r"ppt/slides/slide\d+\.xml", name))
    explicit_colors = Counter(value.decode("ascii").upper()
                              for value in re.findall(rb'<a:srgbClr val="([0-9A-Fa-f]{6})"', xml))
    palette = master_palette(master_template)
    accepted = set(palette["theme"]["colors"].values()) | set(palette["palette"]["series"])
    accepted |= {value for key, value in palette["palette"].items() if key != "series"}
    accepted |= set(master_color_map(master_template).values())
    off_theme = {color: count for color, count in explicit_colors.items() if color not in accepted}
    warnings = []
    if exact != len(source_hashes):
        warnings.append({"code": "master_parts_changed", "count": len(source_hashes) - exact})
    if out_of_bounds:
        warnings.append({"code": "slide_shape_out_of_bounds", "count": len(out_of_bounds)})
    if fonts:
        warnings.append({"code": "font_override_present", "fonts": dict(fonts)})
    if off_theme:
        warnings.append({"code": "off_theme_colors", "colors": off_theme})
    return {
        "status": "pass" if not warnings else "warning",
        "slide_count": len(prs.slides),
        "layout_usage": dict(Counter(str(slide.slide_layout.name) for slide in prs.slides)),
        "master_parts": {"source": len(source_hashes), "exact": exact},
        "filled_placeholders": filled,
        "out_of_bounds": out_of_bounds,
        "explicit_font_overrides": dict(fonts),
        "explicit_colors": dict(explicit_colors),
        "off_theme_colors": off_theme,
        "warnings": warnings,
    }
