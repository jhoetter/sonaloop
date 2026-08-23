"""Decision-oriented visual builders for evidence-linked presentation plans."""
from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ._deck import PALETTE as _PALETTE

_INK = _PALETTE["ink"]
_MUTED = _PALETTE["muted"]
_ACCENT = _PALETTE["accent"]
_ACCENT_INK = _PALETTE["accentInk"]
_PANEL = _PALETTE["panel"]
_LINE = _PALETTE["line"]
_SURFACE2 = _PALETTE["surface2"]
_ACCENT_WEAK = _PALETTE["accentWeak"]


def _place_image(slide, path, left, top, width, height):
    if not path:
        return None
    try:
        pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top))
        scale = min(Inches(width) / pic.width, Inches(height) / pic.height)
        pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
        pic.left = Inches(left) + (Inches(width) - pic.width) // 2
        pic.top = Inches(top) + (Inches(height) - pic.height) // 2
        return pic
    except Exception:
        return None


def build_decision_dashboard(s, e):
    """A dense executive readout: decision, proof and rationale in one visual frame."""
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Decision")
    decision = s.get("decision") or {}
    left, top, left_w = 0.7, 1.72, 7.12
    e.rrect(slide, left, top, left_w, 2.45, _ACCENT_WEAK, radius=0.05)
    label = str(decision.get("label") or decision.get("status") or "Decision").upper()
    lt = e.text(slide, left + 0.35, top + 0.28, left_w - 0.7, 0.28,
                label, size=9.5, bold=True, color=_ACCENT)
    e.mono_run(lt.text_frame.paragraphs[0].runs[0])
    e.text(slide, left + 0.35, top + 0.72, left_w - 0.7, 0.95,
           str(decision.get("text") or decision.get("title") or s.get("heading") or ""),
           size=28, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if decision.get("detail"):
        e.text(slide, left + 0.35, top + 1.82, left_w - 0.7, 0.38,
               str(decision["detail"]), size=11, color=_MUTED)

    metrics = list(s.get("metrics") or [])[:2]
    metric_x, metric_w, gap = 8.14, W.inches - 8.14 - 0.7, 0.18
    metric_h = (2.45 - gap * max(0, len(metrics) - 1)) / max(1, len(metrics))
    for index, item in enumerate(metrics):
        y = top + index * (metric_h + gap)
        e.rrect(slide, metric_x, y, metric_w, metric_h, _PANEL,
                radius=0.05, line=_ACCENT if index == 0 else _LINE)
        tf = e.box(slide, Inches(metric_x + 0.28), Inches(y + 0.15),
                   Inches(metric_w - 0.56), Inches(metric_h - 0.3))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        e.run(tf.paragraphs[0], str(item.get("value") or ""), size=30, bold=True,
              color=_INK)
        p = tf.add_paragraph(); p.space_before = Pt(3)
        e.run(p, str(item.get("label") or ""), size=10.5, bold=True)
        if item.get("detail") or item.get("sub"):
            pd = tf.add_paragraph(); pd.space_before = Pt(2)
            e.run(pd, str(item.get("detail") or item.get("sub")), size=9.5, color=_MUTED)

    rationale = list(s.get("rationale") or [])[:3]
    if rationale:
        rat_top, rat_gap = 4.48, 0.22
        rat_w = (W.inches - 1.4 - rat_gap * (len(rationale) - 1)) / len(rationale)
        for index, raw_item in enumerate(rationale):
            item = raw_item if isinstance(raw_item, dict) else {"text": str(raw_item)}
            x = 0.7 + index * (rat_w + rat_gap)
            e.rrect(slide, x, rat_top, rat_w, 1.72, _PANEL, radius=0.04, line=_LINE)
            e.rrect(slide, x + 0.25, rat_top + 0.25, 0.05, 0.25,
                    _ACCENT, radius=0.3)
            e.text(slide, x + 0.43, rat_top + 0.18, rat_w - 0.68, 0.35,
                   str(item.get("title") or item.get("label") or ""), size=12, bold=True)
            e.text(slide, x + 0.25, rat_top + 0.68, rat_w - 0.5, 0.72,
                   str(item.get("text") or item.get("detail") or ""),
                   size=10.5, color=_MUTED, anchor=MSO_ANCHOR.TOP)
    e.footer(slide)


def build_revision_mockup(s, e):
    """Show the admitted current stimulus beside an editable, concrete proposed revision."""
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Proposed revision")
    left, top, panel_h = 0.7, 1.72, H.inches - 2.22
    left_w, gap = 5.55, 0.35
    right, right_w = left + left_w + gap, W.inches - (left + left_w + gap) - 0.7

    e.rrect(slide, left, top, left_w, panel_h, _PANEL, radius=0.04, line=_LINE)
    source_label = e.text(slide, left + 0.25, top + 0.2, left_w - 0.5, 0.28,
                          str(s.get("source_label") or "Current").upper(),
                          size=9, bold=True, color=_MUTED)
    e.mono_run(source_label.text_frame.paragraphs[0].runs[0])
    pic = _place_image(slide, s.get("image"), left + 0.25, top + 0.62,
                       left_w - 0.5, panel_h - 0.88)
    if pic is None:
        e.rrect(slide, left + 0.25, top + 0.62, left_w - 0.5, panel_h - 0.88,
                _SURFACE2, radius=0.02, line=_LINE)

    proposal = s.get("proposal") or {}
    e.rrect(slide, right, top, right_w, panel_h, _ACCENT_WEAK,
            radius=0.04, line=_ACCENT)
    proposal_label = e.text(slide, right + 0.3, top + 0.2, right_w - 0.6, 0.28,
                            str(s.get("proposal_label") or "Proposed").upper(),
                            size=9, bold=True, color=_ACCENT)
    e.mono_run(proposal_label.text_frame.paragraphs[0].runs[0])
    card_x, card_y = right + 0.3, top + 0.64
    card_w, card_h = right_w - 0.6, 3.55
    e.rrect(slide, card_x, card_y, card_w, card_h, _PANEL, radius=0.04, line=_LINE)
    if proposal.get("eyebrow"):
        eyebrow = e.text(slide, card_x + 0.34, card_y + 0.28, card_w - 0.68, 0.25,
                         str(proposal["eyebrow"]).upper(), size=8.5,
                         bold=True, color=_ACCENT)
        e.mono_run(eyebrow.text_frame.paragraphs[0].runs[0])
    headline_y = card_y + (0.68 if proposal.get("eyebrow") else 0.34)
    e.text(slide, card_x + 0.34, headline_y, card_w - 0.68, 0.62,
           str(proposal.get("headline") or proposal.get("title") or ""),
           size=18, bold=True, anchor=MSO_ANCHOR.TOP)
    e.text(slide, card_x + 0.34, headline_y + 0.74, card_w - 0.68, 1.05,
           str(proposal.get("body") or proposal.get("text") or ""),
           size=11, color=_MUTED, anchor=MSO_ANCHOR.TOP)
    primary = str(proposal.get("primary_cta") or "")
    secondary = str(proposal.get("secondary_cta") or "")
    button_y = card_y + card_h - 0.74
    if primary:
        e.rrect(slide, card_x + 0.34, button_y, 2.15, 0.48, _ACCENT, radius=0.08)
        e.text(slide, card_x + 0.34, button_y + 0.05, 2.15, 0.3, primary,
               size=9.5, bold=True,
               color=_INK if e.master_mode else _ACCENT_INK,
               align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)
    if secondary:
        e.rrect(slide, card_x + 2.65, button_y, 2.15, 0.48, _PANEL,
                radius=0.08, line=_LINE)
        e.text(slide, card_x + 2.65, button_y + 0.05, 2.15, 0.3, secondary,
               size=9.5, bold=True, color=_INK, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)
    why = [str(item.get("text") or item.get("title") or item)
           if isinstance(item, dict) else str(item) for item in (s.get("why") or [])]
    if why:
        e.text(slide, right + 0.3, top + 4.48, right_w - 0.6, 0.76,
               "  ·  ".join(value for value in why[:3] if value),
               size=9.5, color=_MUTED, anchor=MSO_ANCHOR.TOP)
    e.footer(slide)
