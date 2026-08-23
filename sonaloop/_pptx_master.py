"""PowerPoint master inspection, semantic layout selection and theme extraction.

Customer masters are the authority for presentation geometry and visual identity.  Layout
names and placeholder types are the portable signals shared by PowerPoint, Keynote and
LibreOffice; theme XML plus the artwork actually used by layouts supplies the palette.
"""
from __future__ import annotations

from collections import Counter, defaultdict
import colorsys
import io
import re
from typing import Any, Iterable


_BLANK_NAMES = {
    "blank", "blank slide", "blanko", "blanko folie", "leer", "leere folie",
    "vide", "en blanco",
}


def _name(value: Any) -> str:
    return re.sub(r"[^\w]+", " ", str(value or "").casefold()).strip()


def layout_role(name: str) -> str:
    """Infer a small, platform-neutral role from a presentation layout name."""
    value = _name(name)
    if value in _BLANK_NAMES:
        return "blank"
    if any(token in value for token in (
            "closing", "thank", "questions", "abschluss", "schluss", "danke", "fragen")):
        return "closing"
    if any(token in value for token in (
            "agenda", "contents", "table of contents", "inhaltsverzeichnis", "übersicht")):
        return "agenda"
    if any(token in value for token in (
            "section", "chapter", "divider", "kapitel", "abschnitt", "trennfolie")):
        return "section"
    if (value in {"title slide", "title", "titelfolie", "deckblatt", "cover"}
            or any(token in value for token in (
                "cover", "title slide", "titelfolie", "deckblatt"))):
        return "cover"
    if any(token in value for token in (
            "two content", "two column", "comparison", "zwei inhalte", "zwei spalten",
            "vergleich")):
        return "two_column"
    if any(token in value for token in (
            "image", "picture", "bild", "photo", "foto")):
        return "image"
    if any(token in value for token in (
            "content", "inhalt", "body", "text", "chart", "diagram",
            "title only", "nur titel")):
        return "content"
    return "other"


def _placeholder_names(layout) -> set[str]:
    return {
        str(getattr(ph.placeholder_format.type, "name", ph.placeholder_format.type)).casefold()
        for ph in layout.placeholders
    }


def semantic_layout_role(layout) -> str:
    """Classify a layout by its name, then by portable placeholder structure.

    Agencies routinely rename every layout (sometimes to opaque identifiers).  Placeholder
    structure is the safe fallback: it keeps those masters usable without pretending we can
    distinguish an agenda from a regular content slide when the template itself does not say.
    """
    named = layout_role(getattr(layout, "name", ""))
    if named != "other":
        return named
    placeholders = _placeholder_names(layout)
    body_count = sum(
        1 for ph in layout.placeholders
        if str(getattr(ph.placeholder_format.type, "name", ph.placeholder_format.type)).casefold()
        in {"body", "object"}
    )
    if not placeholders:
        return "blank"
    if "center_title" in placeholders and "subtitle" in placeholders:
        return "cover"
    if "picture" in placeholders:
        return "image"
    if placeholders & {"title", "center_title"} and body_count >= 2:
        return "two_column"
    if placeholders & {"title", "center_title"} and placeholders & {
            "body", "object", "subtitle"}:
        return "content"
    if placeholders & {"title", "center_title"}:
        return "content"
    return "other"


def blank_layout(prs):
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("PowerPoint master contains no slide layouts")
    for layout in layouts:
        if semantic_layout_role(layout) == "blank":
            return layout
    return min(layouts, key=lambda layout: len(layout.placeholders))


def _wanted_roles(kind: str) -> tuple[str, ...]:
    if kind in {"cover", "title"}:
        return ("cover",)
    if kind == "agenda":
        return ("agenda", "content")
    if kind in {"section", "canvas-section"}:
        return ("section", "cover")
    if kind == "closing":
        return ("closing", "section", "cover")
    if kind == "image":
        return ("image", "content")
    if kind in {"comparison", "charts", "stimulus_comparison", "revision_mockup"}:
        return ("two_column", "content")
    if kind in {"persona_grid", "persona_detail", "preference_shift", "annotated_screen",
                "decision_dashboard"}:
        return ("content", "two_column", "image")
    return ("content", "two_column", "image")


def _kind_tokens(kind: str) -> tuple[str, ...]:
    return {
        "cover": ("cover", "title slide", "titelfolie", "deckblatt"),
        "title": ("cover", "title slide", "titelfolie", "deckblatt"),
        "agenda": ("agenda", "contents", "inhaltsverzeichnis", "übersicht"),
        "section": ("section", "kapitel", "abschnitt", "divider"),
        "canvas-section": ("section", "kapitel", "abschnitt", "divider"),
        "closing": ("closing", "abschluss", "schluss", "thank", "danke"),
        "image": ("image", "picture", "bild", "photo", "foto"),
        "comparison": ("two", "zwei", "comparison", "vergleich"),
        "charts": ("two", "zwei", "chart", "diagram"),
    }.get(kind, ("content", "inhalt", "body", "text"))


def layout_for_slide(prs, slide: dict, fallback=None):
    """Choose the closest semantic layout, with placeholder fitness as a tie-breaker."""
    kind = str(slide.get("kind") or "content").casefold()
    wanted = _wanted_roles(kind)
    tokens = _kind_tokens(kind)
    layouts = list(prs.slide_layouts)
    candidates = [layout for layout in layouts
                  if semantic_layout_role(layout) in wanted]
    if not candidates:
        return fallback or blank_layout(prs)

    def score(layout) -> tuple[int, int, int, int]:
        role = semantic_layout_role(layout)
        value = _name(getattr(layout, "name", ""))
        placeholders = _placeholder_names(layout)
        has_title = bool(placeholders & {"title", "center_title"})
        has_body = bool(placeholders & {"body", "object", "subtitle"})
        has_picture = "picture" in placeholders
        fitness = (
            (2 if kind in {"cover", "title", "section", "closing"} and has_title else 0)
            + (2 if kind == "agenda" and has_body else 0)
            + (3 if kind == "image" and has_picture else 0)
            + (1 if kind not in {"cover", "title", "section", "closing"} and has_body else 0)
        )
        return (
            wanted.index(role),
            0 if any(token in value for token in tokens) else 1,
            -fitness,
            len(value),
        )

    return min(candidates, key=score)


def _theme_part(prs):
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    for master in prs.slide_masters:
        try:
            return master.part.part_related_by(RT.THEME)
        except KeyError:
            continue
    return None


def _theme_profile(prs) -> dict[str, Any]:
    from xml.etree import ElementTree as ET

    part = _theme_part(prs)
    if part is None:
        return {"name": "", "colors": {}, "fonts": {}}
    root = ET.fromstring(part.blob)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    scheme = root.find(".//a:clrScheme", ns)
    colors: dict[str, str] = {}
    if scheme is not None:
        for item in list(scheme):
            child = next(iter(item), None)
            value = ((child.attrib.get("lastClr") or child.attrib.get("val"))
                     if child is not None else None)
            if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
                colors[item.tag.rsplit("}", 1)[-1]] = value.upper()
    font_scheme = root.find(".//a:fontScheme", ns)
    fonts: dict[str, str] = {}
    if font_scheme is not None:
        for role, tag in (("major", "majorFont"), ("minor", "minorFont")):
            latin = font_scheme.find(f"a:{tag}/a:latin", ns)
            if latin is not None and latin.attrib.get("typeface"):
                fonts[role] = latin.attrib["typeface"]
    return {
        "name": str(scheme.attrib.get("name") or "") if scheme is not None else "",
        "colors": colors,
        "fonts": fonts,
    }


def _hex_rgb(value: str) -> tuple[float, float, float]:
    value = value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) / 255 for i in (0, 2, 4))


def _mix(a: str, b: str, amount: float) -> str:
    aa, bb = _hex_rgb(a), _hex_rgb(b)
    return "".join(f"{round((aa[i] + (bb[i] - aa[i]) * amount) * 255):02X}"
                   for i in range(3))


def _colorfulness(value: str) -> tuple[float, float, float]:
    h, l, s = colorsys.rgb_to_hls(*_hex_rgb(value))
    return h, l, s


_THEME_NAME_TO_KEY = {
    "DARK_1": "dk1", "TEXT_1": "dk1", "DARK_2": "dk2", "TEXT_2": "dk2",
    "LIGHT_1": "lt1", "BACKGROUND_1": "lt1",
    "LIGHT_2": "lt2", "BACKGROUND_2": "lt2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
}


def _walk_shapes(shapes) -> Iterable[Any]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _intersection_area(shape, width: int, height: int) -> float:
    left = max(0, int(shape.left)); top = max(0, int(shape.top))
    right = min(width, int(shape.left + shape.width))
    bottom = min(height, int(shape.top + shape.height))
    return float(max(0, right - left) * max(0, bottom - top))


def _dominant_layout_color(prs, colors: dict[str, str]) -> str | None:
    weights: defaultdict[str, float] = defaultdict(float)
    for layout in prs.slide_layouts:
        for shape in _walk_shapes(layout.shapes):
            if shape.is_placeholder:
                continue
            area = _intersection_area(shape, int(prs.slide_width), int(prs.slide_height))
            if not area:
                continue
            try:
                fill = shape.fill
                if fill.type is None:
                    continue
                try:
                    direct = fill.fore_color.rgb
                except AttributeError:
                    direct = None
                value = str(direct).upper() if direct is not None else ""
                if not value:
                    theme_name = getattr(fill.fore_color.theme_color, "name", "")
                    value = colors.get(_THEME_NAME_TO_KEY.get(theme_name, ""), "")
                if value:
                    _hue, lightness, saturation = _colorfulness(value)
                    if saturation >= 0.18 and 0.08 < lightness < 0.96:
                        weights[value] += area
            except Exception:
                continue
    return max(weights, key=weights.get) if weights else None


def _closest_hue(candidates: list[str], target: float, fallback: str) -> str:
    viable = []
    for color in candidates:
        hue, lightness, saturation = _colorfulness(color)
        if saturation >= 0.18 and 0.08 < lightness < 0.94:
            distance = min(abs(hue - target), 1 - abs(hue - target))
            viable.append((distance, -saturation, color))
    return min(viable)[2] if viable else fallback


def master_palette(data: bytes) -> dict[str, Any]:
    """Derive a semantic deck palette from the customer's theme and used artwork."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    theme = _theme_profile(prs)
    colors = theme["colors"]
    bg = colors.get("lt1", "FFFFFF")
    ink = colors.get("dk1", "1A1815")
    muted = colors.get("dk2") or _mix(ink, bg, 0.42)
    dominant = _dominant_layout_color(prs, colors)
    accent = dominant or colors.get("accent1") or colors.get("lt2") or "5B5FEF"
    candidates = []
    for value in [accent, *(colors.get(f"accent{i}", "") for i in range(1, 7)),
                  colors.get("lt2", "")]:
        if value and value not in candidates:
            candidates.append(value)
    while len(candidates) < 7:
        candidates.append(accent)
    luminance = sum(channel * weight for channel, weight in zip(
        _hex_rgb(accent), (0.2126, 0.7152, 0.0722)))
    palette = {
        "bg": bg, "panel": bg, "surface2": _mix(bg, ink, 0.055),
        "line": _mix(bg, ink, 0.18), "ink": ink, "muted": muted,
        "faint": _mix(ink, bg, 0.55), "accent": accent,
        "accentInk": "1A1815" if luminance > 0.58 else "FFFFFF",
        "accentWeak": _mix(bg, accent, 0.20),
        "green": _closest_hue(candidates, 0.36, accent),
        "amber": _closest_hue(candidates, 0.12, accent),
        "red": _closest_hue(candidates, 0.98, accent),
        "blue": _closest_hue(candidates, 0.58, accent),
        "violet": _closest_hue(candidates, 0.77, accent),
        "skep": _closest_hue(candidates, 0.07, accent),
    }
    # Keep chart-series and semantic status colors coherent.  Several source colors are
    # intentionally shared by both roles, so divergent targets would make a raw RGB map
    # ambiguous and could turn e.g. a success green into an unrelated brand blue.
    palette["series"] = [palette[key] for key in (
        "accent", "violet", "blue", "green", "amber", "red", "skep",
    )]
    return {"theme": theme, "palette": palette}


def master_color_map(data: bytes) -> dict[str, str]:
    """Map Sonaloop's semantic source colors to colors derived from the master."""
    from ._deck import PALETTE

    palette = master_palette(data)["palette"]
    mapping: dict[str, str] = {}
    # setdefault is deliberate: the stock palette uses white for both panels and text on
    # accents.  Geometry covers far more pixels, so panel/background ownership wins; native
    # master placeholders supply their own contrasting text colors where that role matters.
    for key in (
            "panel", "bg", "surface2", "line", "ink", "muted", "faint", "accent",
            "accentWeak", "green", "amber", "red", "blue", "violet", "skep",
            "accentInk"):
        if key in PALETTE:
            mapping.setdefault(str(PALETTE[key]).upper(), palette[key])
    for source, target in zip(PALETTE["series"], palette["series"]):
        mapping.setdefault(str(source).upper(), target)
    mapping.update({
        "E7F3EC": _mix(palette["bg"], palette["green"], 0.14),
        "F6ECDD": _mix(palette["bg"], palette["amber"], 0.14),
        "E4EEF7": _mix(palette["bg"], palette["blue"], 0.14),
        "F7E6E9": _mix(palette["bg"], palette["red"], 0.14),
    })
    return mapping


def inspect_master_template(data: bytes) -> dict[str, Any]:
    """Return a safe layout/theme profile without retaining customer slide content."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    layouts = []
    roles: Counter[str] = Counter()
    for index, layout in enumerate(prs.slide_layouts):
        role = semantic_layout_role(layout)
        roles[role] += 1
        placeholders = []
        for placeholder in layout.placeholders:
            raw = placeholder.placeholder_format.type
            placeholders.append({
                "type": str(getattr(raw, "name", raw)).casefold(),
                "x": round(placeholder.left / 914400, 3),
                "y": round(placeholder.top / 914400, 3),
                "width": round(placeholder.width / 914400, 3),
                "height": round(placeholder.height / 914400, 3),
            })
        layouts.append({
            "index": index,
            "name": str(getattr(layout, "name", "") or f"Layout {index + 1}"),
            "role": role,
            "placeholder_count": len(placeholders),
            "placeholder_types": [item["type"] for item in placeholders],
            "placeholders": placeholders,
        })
    theme = master_palette(data)
    available = set(roles)
    coverage = {
        "cover": bool(available & {"cover"}),
        "content": bool(available & {"content", "two_column", "image"}),
        "agenda": bool(available & {"agenda", "content"}),
        "section": bool(available & {"section", "cover"}),
        "closing": bool(available & {"closing", "section", "cover"}),
        "image": bool(available & {"image", "content"}),
    }
    warnings = []
    if not coverage["cover"]:
        warnings.append({"code": "missing_cover_layout"})
    if not coverage["content"]:
        warnings.append({"code": "missing_content_layout"})
    if not any(item["placeholder_count"] for item in layouts):
        warnings.append({"code": "no_placeholders"})
    return {
        "profile_version": 2,
        "layout_count": len(layouts),
        "layouts": layouts,
        "role_counts": dict(sorted(roles.items())),
        "theme": theme["theme"],
        "semantic_palette": theme["palette"],
        "compatibility": {
            "status": "ready" if not warnings else "limited",
            "coverage": coverage,
            "warnings": warnings,
        },
    }
