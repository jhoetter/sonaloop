"""Slide builders for the PPTX renderer — one `build_<kind>(spec, eng)` per deck layout.

Split out of sonaloop/_pptx.py to keep both modules under the LOC bar (spec/refactor-plan.md),
behaviour-preserving. `_pptx.render()` constructs an engine (the captured Presentation + the
shared drawing primitives), then dispatches each slide spec to the matching builder here. The
builders paint through `eng.<primitive>(…)` (the SAME primitives the chart painters use via
`_pptx_charts.draw`), so the slide, chart and builder layers can never drift apart. The slide
model + chart shapes are documented in _pptx.py.
"""
from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette — same vendored source as _pptx.py (sonaloop-design deck.data.mjs → _deck.py).
from ._deck import PALETTE as _PALETTE, TONES as _TONES, TYPE as _TYPE

_INK = _PALETTE["ink"]
_MUTED = _PALETTE["muted"]
_FAINT = _PALETTE["faint"]
_ACCENT = _PALETTE["accent"]
_BG = _PALETTE["bg"]
_PANEL = _PALETTE["panel"]
_SERIES = list(_PALETTE["series"])
_LINE = _PALETTE["line"]
_SURFACE2 = _PALETTE["surface2"]
_GREEN, _AMBER, _RED = _PALETTE["green"], _PALETTE["amber"], _PALETTE["red"]
_ACCENT_WEAK = _PALETTE["accentWeak"]
_TS = {k: v.get("size", 13) for k, v in _TYPE.items()}

from ._pptx_charts import _num

# canonical stance terms → tone colour (stance_scale.json roles; compatibility "opposed" kept).
_SENTIMENT = {"support": _GREEN, "conditional": _AMBER, "neutral": _MUTED,
              "skeptical": _AMBER, "oppose": _RED, "opposed": _RED}
_CALLOUT_RGB = {"accent": _ACCENT, "green": _GREEN, "amber": _AMBER}
_flow_sz = {"li": 13, "quote": 15, "h": 15}


# ── title slide ──────────────────────────────────────────────────────────
def build_title(s, e):
    prs, blank, W, H, title = e.prs, e.blank, e.W, e.H, e.title
    slide = prs.slides.add_slide(blank)
    e.bg(slide)
    e.rule(slide, 0.92, 2.0, 1.0)
    tf = e.box(slide, Inches(0.9), Inches(2.15), W - Inches(1.8), Inches(3.4))
    p0 = tf.paragraphs[0]
    eb = e.run(p0, (s.get("eyebrow", "") or "").upper(), size=12, bold=True, color=_ACCENT)
    e.mono_run(eb)
    p1 = tf.add_paragraph(); p1.space_before = Pt(8)
    e.run(p1, s.get("title", title), size=38, bold=True)
    if s.get("subtitle"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(12)
        e.run(p2, s["subtitle"], size=14, color=_MUTED)
    if s.get("lead"):
        p3 = tf.add_paragraph(); p3.space_before = Pt(20)
        e.run(p3, s["lead"], size=17, color=_INK)
    e.footer(slide)


# ── content slide ────────────────────────────────────────────────────────
def build_content(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank)
    e.bg(slide)
    # heading: mono section number + bold title, accent rule beneath
    htf = e.box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
    hp = htf.paragraphs[0]
    if s.get("num"):
        e.mono_run(e.run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT))
    e.run(hp, s.get("heading", ""), size=24, bold=True)
    e.rule(slide, 0.72, 1.34, 0.85)

    has_visual = bool(s.get("chart") or s.get("image"))
    # prose caps at ~10in without a visual (full-frame lines read like terminal output).
    body_w = 6.4 if has_visual else min(W.inches - 1.4, 10.0)
    region_top = 1.62
    region_bot = H.inches - (0.78 if s.get("footnote") else 0.5)
    region_h = region_bot - region_top

    # Segment the blocks: runs of flowing text (p/li/h/quote) interleaved with boxed callouts.
    segs = []; run_blocks = []
    def _flush_run():
        if run_blocks:
            h = sum(e.est_h(b.get("text", ""), body_w - 0.2, _flow_sz.get(b.get("type", "p"), 14)) + 0.10
                    for b in run_blocks)
            segs.append(("text", list(run_blocks), h)); run_blocks.clear()
    for b in (s.get("blocks") or []):
        if b.get("type") == "callout":
            _flush_run()
            bh = max(0.66, 0.4 + (0.26 if b.get("label") else 0) + e.est_h(b.get("text", ""), body_w - 0.5, 13))
            segs.append(("callout", b, bh))
        else:
            run_blocks.append(b)
    _flush_run()

    gap = 0.16
    total = sum(h for _, _, h in segs) + gap * max(0, len(segs) - 1)
    y = region_top + (max(0.0, (region_h - total) / 2) if total < region_h else 0.0)
    for kind_, payload, h in segs:
        if kind_ == "callout":
            e.callout_box(slide, 0.7, y, body_w, payload)
        else:
            tf = e.box(slide, Inches(0.7), Inches(y), Inches(body_w), Inches(h + 0.4))
            first = True
            for b in payload:
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                bt = b.get("type", "p"); txt = b.get("text", "")
                if bt == "li":
                    p.space_after = Pt(5)
                    e.run(p, "•  ", size=13, bold=True, color=_ACCENT); e.run(p, txt, size=13, color=_INK)
                elif bt == "quote":
                    p.space_before = Pt(4); p.space_after = Pt(6); e.run(p, txt, size=15, italic=True, color=_MUTED)
                elif bt == "h":
                    p.space_before = Pt(8); p.space_after = Pt(2); e.run(p, txt, size=15, bold=True, color=_INK)
                else:
                    p.space_after = Pt(6); e.run(p, txt, size=14, color=_INK)
        y += h + gap

    rx = 7.35; rw = W.inches - rx - 0.7
    if s.get("chart"):
        e.chart(slide, s["chart"], Inches(rx), Inches(region_top), Inches(rw), Inches(region_h))
    elif s.get("image"):
        try:
            pic = slide.shapes.add_picture(s["image"], Inches(rx), Inches(region_top))
            sc = min(Inches(rw) / pic.width, Inches(region_h) / pic.height)
            pic.width = int(pic.width * sc); pic.height = int(pic.height * sc)
            pic.left = Inches(rx) + (Inches(rw) - pic.width) // 2
            pic.top = Inches(region_top) + (Inches(region_h) - pic.height) // 2
            pic.line.color.rgb = e.rgb(_LINE); pic.line.width = Pt(0.75)
        except Exception:
            pass
    if s.get("footnote"):
        ftf = e.box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.5))
        e.run(ftf.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
    e.footer(slide)


# ── image slide (prototype screenshots / images / avatars) — fitted + centred ──
def build_image(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank)
    e.bg(slide)
    htf = e.box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
    hp = htf.paragraphs[0]
    if s.get("num"):
        e.mono_run(e.run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT))
    e.run(hp, s.get("heading", ""), size=24, bold=True)
    e.rule(slide, 0.72, 1.34, 0.85)
    L, T = 0.7, 1.7
    maxw, maxh = W.inches - 1.4, H.inches - T - 0.95
    placed = False
    if s.get("image"):
        try:
            pic = slide.shapes.add_picture(s["image"], Inches(L), Inches(T))
            scale = min(Inches(maxw) / pic.width, Inches(maxh) / pic.height)
            pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
            pic.left = Inches(L) + (Inches(maxw) - pic.width) // 2
            pic.top = Inches(T) + (Inches(maxh) - pic.height) // 2
            pic.line.color.rgb = e.rgb(_LINE); pic.line.width = Pt(0.75)
            if s.get("caption"):
                cap_t = (pic.top + pic.height) / 914400 + 0.06
                e.text(slide, L, cap_t, maxw, 0.3, s["caption"], size=10, color=_MUTED, align=PP_ALIGN.CENTER)
            placed = True
        except Exception:
            pass
    if not placed:
        # missing/unloadable file → a quiet placeholder panel (the master-template docs case)
        ph = maxh - (0.35 if s.get("caption") else 0)
        e.rrect(slide, L, T, maxw, ph, _SURFACE2, radius=0.03, line=_LINE)
        e.text(slide, L, T + ph / 2 - 0.2, maxw, 0.4, "image — fitted & centred",
               size=12, color=_FAINT, align=PP_ALIGN.CENTER)
        if s.get("caption"):
            e.text(slide, L, T + ph + 0.08, maxw, 0.3, s["caption"], size=10, color=_MUTED, align=PP_ALIGN.CENTER)
    e.footer(slide)


def build_cover(s, e):
    prs, blank, W, H, title = e.prs, e.blank, e.W, e.H, e.title
    slide = prs.slides.add_slide(blank); e.bg(slide)
    canvas = e.da.CANVASES.get(s.get("canvas") or "")
    text_w = W.inches - (4.2 + 1.3 if canvas else 1.8)
    if canvas:
        e.pic_cover(slide, canvas, W.inches - 4.2, 0, 4.2, H.inches)
    if s.get("logo"):
        e.logo_row(slide, 0.9, 0.55)
    e.rule(slide, 0.92, 2.0, 1.0)
    tf = e.box(slide, Inches(0.9), Inches(2.15), Inches(text_w), Inches(3.4))
    e.mono_run(e.run(tf.paragraphs[0], (s.get("eyebrow", "") or "").upper(),
                     size=_TS["eyebrow"], bold=True, color=_ACCENT))
    p1 = tf.add_paragraph(); p1.space_before = Pt(10)
    e.run(p1, s.get("title", title), size=_TS["display"], bold=True)
    if s.get("subtitle"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(12)
        e.run(p2, s["subtitle"], size=16, color=_MUTED)
    if s.get("meta"):
        mt = e.box(slide, Inches(0.9), H - Inches(1.0), Inches(text_w - 2.6), Inches(0.4))
        e.mono_run(e.run(mt.paragraphs[0], s["meta"], size=11, color=_FAINT))
    if s.get("date"):
        dt = e.box(slide, Inches(0.9 + text_w - 2.5), H - Inches(1.0), Inches(2.5), Inches(0.4))
        pd = dt.paragraphs[0]; pd.alignment = PP_ALIGN.RIGHT
        e.run(pd, s["date"], size=11, color=_MUTED)


def build_agenda(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Contents")
    for i, item in enumerate(s.get("items") or []):
        y = 1.95 + i * 0.62
        nt = e.text(slide, 0.9, y, 0.55, 0.5, str(i + 1).zfill(2), size=14, bold=True, color=_ACCENT)
        e.mono_run(nt.text_frame.paragraphs[0].runs[0])
        e.text(slide, 1.5, y, W.inches - 2.6, 0.5, str(item), size=16)
    e.footer(slide)


def build_section(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    bt = e.text(slide, 0.85, 1.0, W.inches - 1.8, 2.6, s.get("num", ""),
                size=_TS["bignum"], bold=True, color=_ACCENT_WEAK, anchor=MSO_ANCHOR.TOP)
    if s.get("num"):
        e.mono_run(bt.text_frame.paragraphs[0].runs[0])
    e.rule(slide, 0.92, 3.95, 0.85)
    tf = e.box(slide, Inches(0.9), Inches(4.15), W - Inches(1.8), Inches(2.2))
    e.run(tf.paragraphs[0], s.get("title", ""), size=32, bold=True)
    if s.get("subtitle"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(10)
        e.run(p2, s["subtitle"], size=_TS["subtitle"], color=_MUTED)
    e.footer(slide)


def build_canvas_section(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    canvas = e.da.CANVASES.get(s.get("canvas") or "")
    if canvas:
        e.pic_cover(slide, canvas, 0, 0, W.inches, H.inches)
    y0 = H.inches - 2.65
    e.rrect(slide, 0.9, y0, 6.4, 1.75, _PANEL, radius=0.09)
    # the card content flows (mirrors the preview): num is optional, title/subtitle move up
    ty = y0 + 0.24
    if s.get("num"):
        nt = e.text(slide, 1.24, y0 + 0.22, 5.7, 0.28, s["num"], size=13, bold=True,
                    color=_ACCENT, anchor=MSO_ANCHOR.TOP)
        e.mono_run(nt.text_frame.paragraphs[0].runs[0])
        ty = y0 + 0.52
    e.text(slide, 1.24, ty, 5.7, 0.55, s.get("title", ""), size=_TS["title"],
           bold=True, anchor=MSO_ANCHOR.TOP)
    if s.get("subtitle"):
        e.text(slide, 1.24, ty + 0.6, 5.7, 0.5, s["subtitle"], size=12, color=_MUTED,
               anchor=MSO_ANCHOR.TOP)


def build_pillars(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    items = s.get("items") or []
    n = max(len(items), 1)
    gap = 0.3
    cw = (W.inches - 1.4 - gap * (n - 1)) / n
    cy = 2.0
    for i, it in enumerate(items):
        x = 0.7 + i * (cw + gap)
        e.rrect(slide, x, cy, 0.72, 0.72, _ACCENT_WEAK, radius=0.19)
        icon = e.da.ICONS.get(it.get("icon") or "", {}).get("accent")
        if icon:
            e.pic(slide, icon, x + 0.13, cy + 0.13, 0.46, 0.46)
        e.text(slide, x, cy + 0.9, cw, 0.4, it.get("title", ""), size=14, bold=True,
               anchor=MSO_ANCHOR.TOP)
        e.text(slide, x, cy + 1.32, cw, H.inches - cy - 2.3, it.get("text", ""), size=11.5,
               color=_MUTED, anchor=MSO_ANCHOR.TOP)
    e.footer(slide)


def build_summary(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Executive summary")
    items = s.get("items") or []
    for it, (cx, cy, cw, ch) in zip(items, e.grid_cells(items)):
        e.rrect(slide, cx, cy, cw, ch, _PANEL, radius=0.05, line=_LINE)
        if it.get("icon") and e.icon_chip(slide, cx + 0.28, cy + 0.26, 0.52, it["icon"]):
            tx0 = cx + 0.96
        else:
            e.rrect(slide, cx + 0.26, cy + 0.27, 0.05, 0.21, _ACCENT, radius=0.3); tx0 = cx + 0.42
        # ONE flowing frame (title → body → meta): a wrapping title can never overlap the body
        tf = e.box(slide, Inches(tx0), Inches(cy + 0.18), Inches(cx + cw - 0.26 - tx0), Inches(ch - 0.36))
        e.run(tf.paragraphs[0], it.get("title", ""), size=15, bold=True)
        if it.get("text"):
            pb = tf.add_paragraph(); pb.space_before = Pt(6)
            e.run(pb, it["text"], size=12, color=_MUTED)
        if it.get("meta"):     # quiet card meta (e.g. effort·value scores), mono + faint
            pm = tf.add_paragraph(); pm.space_before = Pt(6)
            e.mono_run(e.run(pm, str(it["meta"]), size=9, color=_FAINT))
    e.footer(slide)


def build_insight(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    tone = _TONES.get(s.get("tone") or s.get("kind") or "insight") or _TONES["insight"]
    tc = _PALETTE.get(tone.get("color", "accent"), _ACCENT)
    et = e.box(slide, Inches(0.7), Inches(0.55), W - Inches(1.4), Inches(0.4))
    ep = et.paragraphs[0]
    e.mono_run(e.run(ep, str(s.get("eyebrow") or tone.get("label", "Insight")).upper(),
                     size=_TS["eyebrow"], bold=True, color=tc))
    if s.get("num"):
        e.mono_run(e.run(ep, "  ·  " + s["num"], size=_TS["eyebrow"], bold=True, color=_FAINT))
    has_chart = bool(s.get("chart"))
    body_w = 6.9 if has_chart else W.inches - 1.4
    stmt_size, sup_size = 30, 15
    # Content region between the eyebrow and the footnote band; balance the block in it.
    region_top = 1.4
    region_bot = H.inches - (0.95 if (s.get("meta") or s.get("footnote")) else 0.6)
    region_h = region_bot - region_top
    # full-height accent bar (always reads intentional, no matter the block height)
    e.rrect(slide, 0.7, region_top, 0.055, region_h, tc, radius=0.3)
    support = [str(t) for t in (s.get("support") or [])]
    est = e.est_h(s.get("statement", ""), body_w - 0.3, stmt_size) + 0.18
    est += sum(e.est_h(t, body_w - 0.5, sup_size) + 0.11 for t in support)
    anchor = MSO_ANCHOR.MIDDLE if est < region_h - 0.2 else MSO_ANCHOR.TOP
    cf = e.box(slide, Inches(0.98), Inches(region_top), Inches(body_w - 0.28), Inches(region_h))
    cf.vertical_anchor = anchor
    e.run(cf.paragraphs[0], s.get("statement", ""), size=stmt_size, bold=True)
    for i, t in enumerate(support):
        sp = cf.add_paragraph()
        sp.space_before = Pt(16 if i == 0 else 7); sp.space_after = Pt(0)
        e.run(sp, "•  ", size=sup_size, bold=True, color=tc)
        e.run(sp, t, size=sup_size)
    if has_chart:
        e.chart(slide, s["chart"], Inches(7.55), Inches(1.55), Inches(5.05), Inches(4.5))
    if s.get("meta"):
        mt = e.box(slide, Inches(0.7), H - Inches(0.85), Inches(7.0), Inches(0.3))
        e.mono_run(e.run(mt.paragraphs[0], s["meta"], size=11, bold=True, color=tc))
    if s.get("footnote"):
        ft = e.box(slide, Inches(0.7), H - Inches(0.6), W - Inches(1.4), Inches(0.3))
        e.run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
    e.footer(slide)


def build_quote(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.text(slide, 1.35, 1.05, 1.6, 1.4, "“", size=120, bold=True, color=_ACCENT_WEAK,
           anchor=MSO_ANCHOR.TOP)
    # large, vertically-centred quote — the emotional anchor of the chapter
    e.text(slide, 1.9, 1.65, W.inches - 3.4, 3.75, s.get("text", ""), size=30,
           anchor=MSO_ANCHOR.MIDDLE)
    e.initials_chip(slide, 1.9, 5.95, 0.36, s.get("attribution"))
    nt = e.text(slide, 2.42, 5.97, W.inches - 4.3, 0.4, s.get("attribution", ""), size=13, bold=True)
    e.run(nt.text_frame.paragraphs[0], "   " + s.get("role", ""), size=11, color=_MUTED)
    e.footer(slide)


def build_voices(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Voices")
    items = s.get("items") or []
    if len(items) <= 2:                       # 1–2 voices: full-width stacked cards (quote-scale)
        gap = 0.25
        chh = (H.inches - 1.75 - 0.85 - gap) / 2
        cells = [(0.7, 1.75 + i * (chh + gap), W.inches - 1.4, chh) for i in range(len(items))]
    else:
        cells = e.grid_cells(items)
    for it, (cx, cy, cw, ch) in zip(items, cells):
        e.rrect(slide, cx, cy, cw, ch, _PANEL, radius=0.05, line=_LINE)
        e.initials_chip(slide, cx + 0.24, cy + 0.18, 0.3, it.get("name"))
        nt = e.text(slide, cx + 0.62, cy + 0.16, cw - 2.3, 0.34, it.get("name", ""), size=12, bold=True)
        e.run(nt.text_frame.paragraphs[0], "   " + it.get("role", ""), size=10, color=_MUTED)
        sc = _SENTIMENT.get((it.get("sentiment") or "").lower(), _MUTED)
        st = e.text(slide, cx + cw - 1.85, cy + 0.16, 1.6, 0.34,
                    str(it.get("sentiment_label") or it.get("sentiment") or "").upper(),
                    size=9, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        if it.get("sentiment"):
            e.mono_run(st.text_frame.paragraphs[0].runs[0])
        bf = e.box(slide, Inches(cx + 0.24), Inches(cy + 0.62), Inches(cw - 0.48), Inches(ch - 0.8))
        e.run(bf.paragraphs[0], it.get("text", ""), size=12 if len(items) > 2 else 13)
    e.footer(slide)


def build_stats(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    items = s.get("items") or []
    n = max(len(items), 1); gap = 0.28
    tw = (W.inches - 1.4 - gap * (n - 1)) / n
    th = 3.0; ty = 1.75 + (H.inches - 1.75 - 0.6 - th) / 2     # centre the KPI band below the heading
    for i, it in enumerate(items):
        tx = 0.7 + i * (tw + gap)
        e.rrect(slide, tx, ty, tw, th, _PANEL, radius=0.05, line=_LINE)
        if it.get("icon") and e.icon_chip(slide, tx + 0.28, ty + 0.32, 0.66, it["icon"]):
            tf = e.box(slide, Inches(tx + 0.28), Inches(ty + 1.18), Inches(tw - 0.52), Inches(th - 1.4))
            tf.vertical_anchor = MSO_ANCHOR.TOP
        else:
            e.rrect(slide, tx + 0.28, ty + 0.32, 0.05, 0.34, _ACCENT, radius=0.3)
            tf = e.box(slide, Inches(tx + 0.28), Inches(ty + 0.24), Inches(tw - 0.52), Inches(th - 0.48))
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        v = it.get("value")
        e.run(tf.paragraphs[0], v if isinstance(v, str) else _num(v or 0), size=40, bold=True)
        pl = tf.add_paragraph(); pl.space_before = Pt(8)
        e.run(pl, str(it.get("label", "")), size=12, color=_MUTED)
        if it.get("sub"):
            p2 = tf.add_paragraph(); p2.space_before = Pt(3)
            e.run(p2, str(it["sub"]), size=10, color=_FAINT)
    e.footer(slide)


def build_chart(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    if s.get("chart"):
        e.chart(slide, s["chart"], Inches(0.7), Inches(1.8), Inches(W.inches - 1.4), Inches(H.inches - 2.9))
    if s.get("footnote"):
        ft = e.box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
        e.run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
    e.footer(slide)


def build_charts(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    items = s.get("items") or []
    n = max(len(items), 1)
    gap = 0.4
    cw = (W.inches - 1.4 - gap * (n - 1)) / n
    top = 1.75
    ch = H.inches - top - 0.5 - (1.3 if s.get("footnote") else 1.0)
    for i, it in enumerate(items):
        x = 0.7 + i * (cw + gap)
        e.text(slide, x, top, cw, 0.35, it.get("title", ""), size=13, bold=True,
               anchor=MSO_ANCHOR.TOP)
        if it.get("chart"):
            e.chart(slide, it["chart"], Inches(x), Inches(top + 0.5), Inches(cw), Inches(ch))
    if s.get("footnote"):
        ft = e.box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
        e.run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
    e.footer(slide)


def build_table(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    cols = s.get("columns") or []
    rows = [list(r) for r in (s.get("rows") or [])]
    if cols:
        tw = W.inches - 1.4
        nrows = len(rows) + 1
        avail = H.inches - 1.9 - (0.95 if s.get("footnote") else 0.6)
        row_h = max(0.5, min(0.95, avail / nrows))      # grow rows to fill the frame (capped)
        th = row_h * nrows
        gf = slide.shapes.add_table(nrows, len(cols),
                                    Inches(0.7), Inches(1.9), Inches(tw), Inches(th))
        tbl = gf.table
        # kill the theme's banding/header styling; the cells below carry the deck's own
        tbl.first_row = False
        tbl.horz_banding = False
        for i in range(nrows):
            tbl.rows[i].height = Inches(row_h)
        for j, c in enumerate(cols):
            cell = tbl.cell(0, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = e.rgb(_ACCENT_WEAK)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            e.run(cell.text_frame.paragraphs[0], str(c), size=12.5, bold=True)
        for i, row in enumerate(rows, start=1):
            for j in range(len(cols)):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = e.rgb(_BG if i % 2 else _PANEL)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = True
                val = str(row[j]) if j < len(row) else ""
                e.run(cell.text_frame.paragraphs[0], val, size=12.5,
                      bold=(j == 0), color=_INK if j == 0 else _MUTED)
    if s.get("footnote"):
        ft = e.box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
        e.run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
    e.footer(slide)


def build_comparison(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s)
    gap = 0.3
    cw = (W.inches - 1.4 - gap) / 2
    cy = 1.75; ch = H.inches - cy - 0.85
    for j, (col, accent) in enumerate(((s.get("left") or {}, False), (s.get("right") or {}, True))):
        cx = 0.7 + j * (cw + gap)
        e.rrect(slide, cx, cy, cw, ch, _PANEL if accent else _SURFACE2, radius=0.04,
                line=_ACCENT if accent else _LINE)
        tf = e.box(slide, Inches(cx + 0.34), Inches(cy + 0.3), Inches(cw - 0.68), Inches(ch - 0.6))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        e.run(tf.paragraphs[0], col.get("title", ""), size=15, bold=True,
              color=_ACCENT if accent else _MUTED)
        tf.paragraphs[0].space_after = Pt(14)
        for t in col.get("items") or []:
            cp = tf.add_paragraph(); cp.space_after = Pt(9)
            e.run(cp, "•  ", size=13, bold=True, color=_ACCENT if accent else _FAINT)
            e.run(cp, str(t), size=13)
    e.footer(slide)


def build_timeline(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    e.heading_band(slide, s, "Next steps")
    steps = s.get("steps") or []
    n = max(len(steps), 1)
    x0, x1, ly = 1.0, W.inches - 1.0, 3.1
    e.connector(slide, x0, ly, x1, ly, _LINE, width=1.2)
    for i, st in enumerate(steps):
        cx = x0 + (i + 0.5) * ((x1 - x0) / n)
        e.oval(slide, cx - 0.1, ly - 0.1, 0.2, _ACCENT)
        if st.get("label"):
            lt = e.text(slide, cx - 1.3, ly - 0.55, 2.6, 0.3, str(st["label"]).upper(),
                        size=10, bold=True, color=_ACCENT, align=PP_ALIGN.CENTER)
            e.mono_run(lt.text_frame.paragraphs[0].runs[0])
        e.text(slide, cx - 1.3, ly + 0.25, 2.6, 0.5, st.get("title", ""), size=13, bold=True,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        e.text(slide, cx - 1.3, ly + 0.72, 2.6, 1.4, st.get("text", ""), size=11, color=_MUTED,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    e.footer(slide)


def build_closing(s, e):
    prs, blank, W, H = e.prs, e.blank, e.W, e.H
    slide = prs.slides.add_slide(blank); e.bg(slide)
    if s.get("logo"):
        e.logo_row(slide, 0.9, 1.1)
    e.rule(slide, 0.92, 2.35, 1.0)
    tf = e.box(slide, Inches(0.9), Inches(2.55), Inches(10.2), Inches(H.inches - 3.6))
    e.run(tf.paragraphs[0], s.get("title", ""), size=40, bold=True)
    if s.get("text"):
        p1 = tf.add_paragraph(); p1.space_before = Pt(18)
        e.run(p1, s["text"], size=16, color=_INK)
    if s.get("contact"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(20)
        e.run(p2, s["contact"], size=14, bold=True, color=_ACCENT)
    if s.get("meta"):
        mt = e.box(slide, Inches(0.9), H - Inches(1.0), W - Inches(1.8), Inches(0.4))
        e.mono_run(e.run(mt.paragraphs[0], s["meta"], size=11, color=_FAINT))


# kind → builder. _pptx.render() dispatches each spec here (fallback = build_content).
PAINTERS = {
    "title": build_title, "cover": build_cover, "agenda": build_agenda,
    "section": build_section, "canvas-section": build_canvas_section,
    "pillars": build_pillars, "summary": build_summary,
    "insight": build_insight, "recommendation": build_insight, "risk": build_insight,
    "quote": build_quote, "voices": build_voices, "stats": build_stats,
    "chart": build_chart, "charts": build_charts, "table": build_table,
    "comparison": build_comparison, "timeline": build_timeline,
    "closing": build_closing, "image": build_image,
}
