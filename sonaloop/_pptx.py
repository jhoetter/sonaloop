"""Native PowerPoint rendering for reports — turns a neutral slide model into a branded .pptx.

This layer is domain-agnostic (no Sonaloop models, no i18n): the report service assembles a list of
plain slide dicts and hands them here. Charts are NATIVE python-pptx charts (editable in PowerPoint),
coloured from the Sonaloop brand series palette so a deck reads on-brand. Lazy-imported by the caller
so the package degrades gracefully when python-pptx is absent.

Slide model (list of dicts):
  {"kind": "title", "title": str, "subtitle": str, "lead": str}
  {"kind": "content", "heading": str,
     "bullets": [(level:int, text:str)],          # level 0 = paragraph, 1+ = nested bullet
     "chart": {...} | None,                         # see below
     "image": "/abs/path.png" | None,               # a figure image (prototype shot / avatar / asset)
     "footnote": str}
Master-template kinds (the deck taxonomy — single source: sonaloop-design/deck.data.mjs,
vendored as sonaloop/_deck.py; every layout is previewed at #/deck in the design docs and
_deck.SAMPLE_SLIDES carries a placeholder example of each):
  {"kind": "cover", "eyebrow", "title", "subtitle", "meta", "date"}
  {"kind": "agenda", "heading", "items": [str]}
  {"kind": "section", "num", "title", "subtitle"}
  {"kind": "summary", "heading", "items": [{"title", "text"}]}
  {"kind": "insight"|"recommendation"|"risk", "tone", "num", "statement",
     "support": [str], "chart": {...}|None, "meta", "footnote"}
  {"kind": "quote", "text", "attribution", "role"}
  {"kind": "voices", "heading", "items": [{"name", "role", "sentiment", "text"}]}
  {"kind": "stats", "heading", "items": [{"label", "value", "sub"}]}
  {"kind": "chart", "num", "heading", "chart": {...}, "footnote"}
  {"kind": "comparison", "heading", "left": {"title", "items"}, "right": {"title", "items"}}
  {"kind": "timeline", "heading", "steps": [{"label", "title", "text"}]}
  {"kind": "closing", "title", "text", "meta", "contact"}
Chart shapes (one per design-system chart `of`):
  {"type": "bar",  "categories": [str], "values": [num]}
  {"type": "pie",  "categories": [str], "values": [num]}
  {"type": "stacked_bar", "rows": [{"label": str, "segments": [{"label": str, "value": num}]}]}
  {"type": "diverging_bar", "rows": [{"label": str, "positive": num, "negative": num}],
     "positive_label": str, "negative_label": str}
  {"type": "gauge", "items": [{"label": str, "value": num, "max": num}]}
  {"type": "dot_plot", "rows": [{"label": str, "values": [num]}], "min": num, "max": num, "unit": str}
  {"type": "heatmap", "columns": [str], "rows": [{"label": str, "values": [num]}]}
  {"type": "line", "series": [{"label": str, "points": [num]}], "labels": [str], "target": num}
  {"type": "stacked_area", "series": [{"label": str, "points": [num]}], "labels": [str]}
  {"type": "column", "items": [{"label": str, "value": num} | {"label": str, "segments": […]}]}
  {"type": "progress_strip", "items": [{"label": str, "value": num}]}
  {"type": "stats", "items": [{"label": str, "value": num|str, "sub": str}]}
  {"type": "scatter", "points": [{"x":num,"y":num,"label":str}], "x_label": str, "y_label": str}
"""
from __future__ import annotations

import base64
import io
from typing import Any

# Sonaloop brand (light) — from the vendored deck master template (_deck.py, generated out of
# sonaloop-design/deck.data.mjs, which derives them from tokens.data.mjs).
from ._deck import PALETTE as _PALETTE, TONES as _TONES, TYPE as _TYPE

_INK = _PALETTE["ink"]
_MUTED = _PALETTE["muted"]
_FAINT = _PALETTE["faint"]
_ACCENT = _PALETTE["accent"]
_BG = _PALETTE["bg"]
_PANEL = _PALETTE["panel"]
# Series palette (accent · violet · blue · green · amber · red · skep).
_SERIES = list(_PALETTE["series"])
_LINE = _PALETTE["line"]
_SURFACE2 = _PALETTE["surface2"]
_GREEN, _AMBER, _RED = _PALETTE["green"], _PALETTE["amber"], _PALETTE["red"]
_ACCENT_WEAK = _PALETTE["accentWeak"]
# Role-based type sizes (pt) from the master template.
_TS = {k: v.get("size", 13) for k, v in _TYPE.items()}


# Number formatting + the chart painters live in _pptx_charts (split for the LOC bar).
from . import _pptx_charts as _pc
from ._pptx_charts import _num
from ._pptx_master import (
    blank_layout as _blank_layout,
    layout_for_slide as _layout_for_slide,
    master_color_map as _master_color_map,
)


def _empty_template_presentation(data: bytes):
    """Open an uploaded .pptx as a master/layout source and remove sample slides.

    python-pptx keeps the slide masters, themes and layouts when slide instances
    are removed. This gives workspace owners a practical master-file workflow
    without executing macros or copying customer content into the generated deck.
    """
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    sld_ids = prs.slides._sldIdLst  # python-pptx has no public slide-delete API
    for slide_id in list(sld_ids):
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        sld_ids.remove(slide_id)
    return prs


def render(slides: list[dict], *, title: str = "Report",
           master_template: bytes | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from types import SimpleNamespace

    master_mode = master_template is not None
    prs = (_empty_template_presentation(master_template)
           if master_template else Presentation())
    if not master_template:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = "Sonaloop"
    blank = _blank_layout(prs)
    W, H = prs.slide_width, prs.slide_height
    color_map = _master_color_map(master_template) if master_template else {}

    def _color(hexv):
        value = str(hexv or "000000").lstrip("#").upper()
        return color_map.get(value, value)

    rgb = lambda hexv: RGBColor.from_string(_color(hexv))

    def _bg(slide, hexv=_BG):
        if master_mode:
            return
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(hexv)

    def _box(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def _run(p, text, *, size=14, bold=False, color=_INK, italic=False):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if not master_mode:
            r.font.name = "Geist"
        r.font.color.rgb = rgb(color)
        return r

    def _rule(slide, l, t, w):
        bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(3.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(_ACCENT); bar.line.fill.background()
        _noshadow(bar)

    def _footer(slide):
        if master_mode:
            return
        ft = _box(slide, W - Inches(5.0), H - Inches(0.42), Inches(4.3), Inches(0.3))
        p = ft.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        _run(p, title, size=9, color=_FAINT)

    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR

    # ── chart primitives: native, EDITABLE shapes, pixel-matched to the .sl-chart components ──
    def _noshadow(sp):
        # autoshapes/connectors carry a <p:style> with a theme effectRef (drop shadow). The DS shapes
        # are flat — drop the style entirely (explicit fill/line are set on every shape anyway).
        try:
            sp.shadow.inherit = False
            st = sp._element.find(qn("p:style"))
            if st is not None:
                sp._element.remove(st)
        except Exception:
            pass

    def _text(slide, l, t, w, h, text, *, size=11, color=_INK, bold=False,
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, rot=0):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        if rot:
            tb.rotation = rot
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]; p.alignment = align
        _run(p, text, size=size, color=color, bold=bold)
        return tb

    def _rrect(slide, l, t, w, h, color, *, radius=0.5, line=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
        if line:
            sp.line.color.rgb = rgb(line); sp.line.width = Pt(1)
        else:
            sp.line.fill.background()
        _noshadow(sp)
        return sp

    def _connector(slide, x1, y1, x2, y2, color, *, width=0.75, dash=False):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = rgb(color); ln.line.width = Pt(width)
        if dash:
            try:
                ln.line._get_or_add_ln().append(ln.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"}))
            except Exception:
                pass
        _noshadow(ln)
        return ln

    def _dot(slide, cxp, cyp, d, fill, edge, num):
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cxp - d / 2), Inches(cyp - d / 2), Inches(d), Inches(d))
        ov.fill.solid(); ov.fill.fore_color.rgb = rgb(fill)
        ov.line.color.rgb = rgb(edge); ov.line.width = Pt(1.5)
        _noshadow(ov)
        tf = ov.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, str(num), size=9, bold=True, color=edge)
        return ov

    # ── presentation-layout helpers: vertical balance + boxed callouts ───────────
    def _est_h(text, width_in, size, *, lf=1.34):
        """Rough rendered height (inches) of wrapped text — used to vertically balance a slide."""
        if not text:
            return size * lf / 72.0
        cpl = max(8, int(width_in / (size * 0.0072)))      # chars per line at this width/size
        line = 1; cur = 0
        for word in str(text).split():
            wl = len(word) + 1
            if cur + wl > cpl and cur:
                line += 1; cur = wl
            else:
                cur += wl
        return line * size * lf / 72.0

    _CALLOUT_BAR = {"accent": _ACCENT, "green": _GREEN, "amber": _AMBER,
                    "blue": _PALETTE.get("blue", _ACCENT), "red": _RED}
    _CALLOUT_TINT = {"accent": _ACCENT_WEAK, "green": "E7F3EC", "amber": "F6ECDD",
                     "blue": "E4EEF7", "red": "F7E6E9"}

    def _callout_box(slide, x, y, w, block, *, size=13):
        """A boxed callout: tinted rounded panel + colour bar + optional icon + label + body. Returns height."""
        kind = block.get("kind") or "accent"
        bar = _CALLOUT_BAR.get(kind, _ACCENT); tint = _CALLOUT_TINT.get(kind, _ACCENT_WEAK)
        label = (block.get("label") or "").strip(); txt = block.get("text") or ""
        icon_b64 = _da.ICONS.get(block.get("icon") or "", {}).get("accent")
        pad = 0.2; inner_w = w - 0.5 - (0.72 if icon_b64 else 0)
        h = max(0.66, pad * 2 + (0.26 if label else 0) + _est_h(txt, inner_w, size))
        _rrect(slide, x, y, w, h, tint, radius=0.06)
        _rrect(slide, x, y + 0.12, 0.07, h - 0.24, bar, radius=0.5)
        tx = x + 0.3
        if icon_b64:
            isz = min(0.5, h - 0.3); _pic(slide, icon_b64, x + 0.28, y + (h - isz) / 2, isz, isz); tx = x + 0.28 + isz + 0.18
        tf = _box(slide, Inches(tx), Inches(y + pad - 0.04), Inches(x + w - 0.24 - tx), Inches(h - pad))
        if label:
            _mono_run(_run(tf.paragraphs[0], label.upper(), size=10, bold=True, color=bar))
            bp = tf.add_paragraph(); bp.space_before = Pt(3)
            _run(bp, txt, size=size, color=_INK)
        else:
            _run(tf.paragraphs[0], txt, size=size, color=_INK)
        return h

    # The chart painters (bar/pie/…/scatter) live in _pptx_charts; they draw through the
    # SAME shape primitives via this ctx, so slide and chart layers can't drift apart.
    from types import SimpleNamespace
    _ctx = SimpleNamespace(text=_text, rrect=_rrect, connector=_connector, dot=_dot,
                           run=_run, noshadow=_noshadow, rgb=rgb, color=_color,
                           font_name=None if master_mode else "Geist")

    def _chart(slide, ch, x, y, cx, cy):
        _pc.draw(_ctx, slide, ch, x, y, cx, cy)

    # ── brand assets — vendored _deck_assets (icons/logos rasterized + canvases recompressed
    # at design time by sonaloop-design/scripts/gen-deck.mjs; PPTX can't embed SVG). Unknown
    # asset names degrade to the unbranded layout. ─────────────────────────────────────
    from . import _deck_assets as _da

    def _pic(slide, b64, l, t, w, h):
        pic = slide.shapes.add_picture(io.BytesIO(base64.b64decode(b64)),
                                       Inches(l), Inches(t), Inches(w), Inches(h))
        _noshadow(pic)
        return pic

    def _pic_cover(slide, b64, l, t, w, h):
        """Aspect-preserving fill (CSS `background: cover`): crop the source to the box ratio."""
        pic = _pic(slide, b64, l, t, w, h)
        try:
            iw, ih = pic.image.size
            sa, ta = iw / ih, w / h
            if sa > ta:
                pic.crop_left = pic.crop_right = (1 - ta / sa) / 2
            elif sa < ta:
                pic.crop_top = pic.crop_bottom = (1 - sa / ta) / 2
        except Exception:
            pass
        return pic

    def _icon_chip(slide, x, y, size, name, *, bg=_ACCENT_WEAK):
        """A hi-fi icon in a tinted rounded chip (the pillars treatment, reusable on any card).
        Returns True if the icon exists in the embedded deck set, else draws nothing."""
        b64 = _da.ICONS.get(name or "", {}).get("accent")
        if not b64:
            return False
        _rrect(slide, x, y, size, size, bg, radius=0.26)
        pad = size * 0.24
        _pic(slide, b64, x + pad, y + pad, size - 2 * pad, size - 2 * pad)
        return True

    def _logo_row(slide, x, y, mark=0.42):
        """The brand moment: mark + wordmark ("sona" ink · "loop" muted)."""
        if master_mode:
            return
        b64 = _da.LOGOS.get("sonaloop")
        if b64:
            _pic(slide, b64, x, y + 0.04, mark, mark)
        tf = _box(slide, Inches(x + mark + 0.14), Inches(y), Inches(2.6), Inches(mark + 0.08))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _run(p, "sona", size=16, bold=True)
        _run(p, "loop", size=16, bold=True, color=_MUTED)

    # ── master-template painters — geometry mirrors sonaloop-design site/deck.preview.mjs
    # painter-for-painter (single source: deck.data.mjs / vendored _deck.py), so the docs
    # previews at #/deck and the exported deck look the same. ──────────────────────────
    def _mono_run(r):
        if not master_mode:
            r.font.name = "Geist Mono"
        return r

    def _heading_band(slide, s, default=""):
        htf = _box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
        hp = htf.paragraphs[0]
        if s.get("num"):
            _mono_run(_run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT))
        _run(hp, s.get("heading", default) or default, size=_TS["title"], bold=True)
        _rule(slide, 0.72, 1.34, 0.85)

    def _grid_cells(items):
        """2-column card grid under the heading band → [(x, y, w, h), …] in inches."""
        rows = max(1, (len(items) + 1) // 2)
        gap = 0.25
        cw = (W.inches - 1.4 - gap) / 2
        ch = (H.inches - 1.75 - 0.85 - gap * (rows - 1)) / rows
        return [(0.7 + (i % 2) * (cw + gap), 1.75 + (i // 2) * (ch + gap), cw, ch)
                for i in range(len(items))]

    def _oval(slide, l, t, d, fill):
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
        ov.fill.solid(); ov.fill.fore_color.rgb = rgb(fill); ov.line.fill.background()
        _noshadow(ov)
        return ov

    def _initials_chip(slide, l, t, d, name):
        ov = _oval(slide, l, t, d, _ACCENT_WEAK)
        tf = ov.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pq = tf.paragraphs[0]; pq.alignment = PP_ALIGN.CENTER
        ini = "".join(w[0] for w in str(name or "?").split()[:2]).upper() or "?"
        _run(pq, ini, size=10, bold=True, color=_ACCENT)

    # ── engine: the captured Presentation state + every drawing primitive. The slide
    # builders (one build_<kind> per layout) live in _pptx_builders and paint THROUGH this
    # engine — the same primitives the chart painters use via _pptx_charts.draw — so the
    # slide, chart and builder layers can never drift apart. ────────────────────────────
    deck_assets = (_da if not master_mode else
                   SimpleNamespace(CANVASES={}, LOGOS={}, ICONS=_da.ICONS))
    eng = SimpleNamespace(
        prs=prs, blank=blank, W=W, H=H, title=title, rgb=rgb, da=deck_assets, num=_num,
        master_mode=master_mode,
        bg=_bg, box=_box, run=_run, rule=_rule, footer=_footer,
        text=_text, rrect=_rrect, connector=_connector, dot=_dot, noshadow=_noshadow,
        est_h=_est_h, callout_box=_callout_box, chart=_chart, pic=_pic, pic_cover=_pic_cover,
        icon_chip=_icon_chip, logo_row=_logo_row, mono_run=_mono_run, heading_band=_heading_band,
        grid_cells=_grid_cells, oval=_oval, initials_chip=_initials_chip,
    )
    from . import _pptx_builders as _b
    from ._pptx_master_native import apply_master_native_layout
    from ._pptx_notes import apply_speaker_notes
    for s in slides:
        eng.blank = _layout_for_slide(prs, s, blank) if master_mode else blank
        _b.PAINTERS.get(s.get("kind"), _b.build_content)(s, eng)
        if master_mode:
            apply_master_native_layout(
                prs.slides[-1], s, slide_width=int(W), slide_height=int(H))
        apply_speaker_notes(prs.slides[-1], s)


    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except Exception:
        return False
