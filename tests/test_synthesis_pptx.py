"""UX V11 — deck/PPTX export quality (spec/ux-contract.md §9): the synthesis→slides mapping uses
the master template's layout vocabulary (statement verdict, takeaway cards, chart/quote/voices
slides, numbered recommendation cards) and NO slide carries a wall of continuous prose — leads are
clamped renderer-side, markdown artifacts and bare artifact ids never reach a slide."""
import io
import re
import zipfile

from pptx import Presentation
from pptx.util import Inches

from sonaloop import services
from sonaloop.services._synthesis import _SYNTHESIS_EXPORT_LABELS
from sonaloop.services._synthesis_pptx import (
    _analytic_slides, _clamp_prose, _label_segments, _split_card,
)


def test_renderer_uses_uploaded_master_layouts_but_discards_sample_slides():
    import io
    from sonaloop import _pptx

    template = Presentation()
    template.slide_width = Inches(13.333)
    template.slide_height = Inches(7.5)
    template.slide_layouts[0].name = "SHKB Cover"
    template.slide_layouts[2].name = "SHKB Section"
    template.slide_layouts[3].name = "SHKB Content"
    template.slide_layouts[6].name = "SHKB Blank"
    template.slides.add_slide(template.slide_layouts[6])
    template.core_properties.subject = "Customer master"
    source = io.BytesIO()
    template.save(source)

    data = _pptx.render([
        {"kind": "cover", "title": "Generated", "subtitle": "From report",
         "speaker_notes": {"talk_track": "Open with the decision context.",
                           "caveats": ["Directional evidence."]}},
        {"kind": "section", "num": "01", "title": "Findings"},
        {"kind": "content", "heading": "Evidence", "blocks": [{"type": "p", "text": "Useful"}]},
        {"kind": "closing", "title": "Done"},
    ], title="Generated", master_template=source.getvalue())
    rendered = Presentation(io.BytesIO(data))

    assert len(rendered.slides) == 4
    assert rendered.core_properties.subject == "Customer master"
    assert rendered.slide_width == Inches(13.333)
    assert rendered.slide_height == Inches(7.5)
    assert [rendered.slides[index].slide_layout.name for index in range(3)] == [
        "SHKB Cover", "SHKB Section", "SHKB Content",
    ]
    assert all(slide._element.cSld.bg is None for slide in rendered.slides)
    assert not any(shape.shape_type == 13 for shape in rendered.slides[0].shapes)
    assert all(run.font.name is None for slide in rendered.slides
               for shape in slide.shapes if shape.has_text_frame
               for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
    notes = rendered.slides[0].notes_slide.notes_text_frame.text
    assert "TALK TRACK" in notes and "Open with the decision context." in notes
    assert "CAVEATS" in notes and "Directional evidence." in notes


def test_master_profile_exposes_semantic_layout_roles():
    from sonaloop._pptx_master import inspect_master_template

    template = Presentation()
    template.slide_layouts[0].name = "SHKB Cover"
    template.slide_layouts[2].name = "SHKB Section Divider"
    template.slide_layouts[3].name = "SHKB Content"
    source = io.BytesIO()
    template.save(source)

    profile = inspect_master_template(source.getvalue())

    assert profile["layout_count"] == len(template.slide_layouts)
    assert profile["role_counts"]["cover"] >= 1
    assert profile["role_counts"]["section"] >= 1
    assert profile["role_counts"]["content"] >= 1
    assert profile["layouts"][0]["name"] == "SHKB Cover"


def test_opaque_layout_names_are_classified_from_placeholder_structure():
    """A customer's agency-specific naming must not be required for master support."""
    from sonaloop._pptx_master import inspect_master_template, layout_for_slide

    template = Presentation()
    for index, layout in enumerate(template.slide_layouts):
        layout.name = f"ACME-{index:02d}"
    source = io.BytesIO()
    template.save(source)

    profile = inspect_master_template(source.getvalue())
    assert profile["compatibility"]["status"] == "ready"
    assert profile["role_counts"]["cover"] >= 1
    assert profile["role_counts"]["content"] >= 1
    assert profile["role_counts"]["two_column"] >= 1
    assert profile["role_counts"]["image"] >= 1
    assert layout_for_slide(template, {"kind": "cover"}).name == "ACME-00"
    assert layout_for_slide(template, {"kind": "comparison"}).name == "ACME-03"
    assert layout_for_slide(template, {"kind": "image"}).name == "ACME-08"


def test_native_agenda_disables_inherited_auto_numbering():
    from pptx.oxml.ns import qn
    from sonaloop import _pptx

    template = Presentation()
    template.slide_layouts[1].name = "Customer Agenda"
    source = io.BytesIO()
    template.save(source)
    data = _pptx.render([
        {"kind": "agenda", "heading": "Inhalt", "items": ["Ergebnis", "Empfehlungen"]},
    ], master_template=source.getvalue())
    rendered = Presentation(io.BytesIO(data))
    body = next(shape for shape in rendered.slides[0].placeholders
                if shape.placeholder_format.type.name.casefold() in {"body", "object"})

    assert body.text_frame.text.splitlines() == ["01   Ergebnis", "02   Empfehlungen"]
    assert all(paragraph._p.pPr.find(qn("a:buNone")) is not None
               for paragraph in body.text_frame.paragraphs)


def test_master_theme_drives_generated_color_and_passes_package_qa():
    from sonaloop import _pptx
    from sonaloop._deck import PALETTE
    from sonaloop._pptx_master import master_color_map, master_palette
    from sonaloop._pptx_master_native import inspect_rendered_master_deck

    template = Presentation()
    raw = io.BytesIO()
    template.save(raw)
    incoming = zipfile.ZipFile(io.BytesIO(raw.getvalue()))
    branded = io.BytesIO()
    with zipfile.ZipFile(branded, "w") as outgoing:
        for item in incoming.infolist():
            value = incoming.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                value = re.sub(
                    rb'(<a:accent1>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}("/>)',
                    rb'\g<1>13A085\2', value,
                )
            outgoing.writestr(item, value)
    master = branded.getvalue()
    slides = [{"kind": "summary", "heading": "Ergebnis", "items": [
        {"title": "Signal", "text": "Aussage aus dem Report"},
    ]}]
    data = _pptx.render(slides, master_template=master)
    semantic = master_palette(master)["palette"]
    mapping = master_color_map(master)

    assert semantic["accent"] == "13A085"
    assert mapping[PALETTE["accent"]] == "13A085"
    assert mapping[PALETTE["panel"]] == semantic["panel"]  # white panel never becomes dark ink
    assert inspect_rendered_master_deck(data, master)["status"] == "pass"

# A showcase-shaped convergence synthesis: caps-label exec prose, stanced voices with verbatim
# council quotes, scored recommendations and label-led findings — the report shape the owner's
# screenshots came from.
_BIG = (
    "THE WINNING CONCEPT: A lunch window that stays visible in the day and slides audibly instead "
    "of dissolving in silence, with hard walls for fixed commitments. When the window hits a wall, "
    "the system says so honestly instead of sliding an illusion. "
    "WHY IT WINS: It makes the break as loud as its disruptors and productises the only defence "
    "observed in the field. The mid-fi validation was strict and every objection was retested. "
    "SURVEY BACKING (5 answers): The problem is real — 4/5 write off their break before 10 am on "
    "busy days. The visible window holds: 1 for, 3 conditional, 1 skeptical."
)


def _showcase(store):
    council = {
        "id": "c1", "created_at": "2026-06-10T00:00:00+00:00", "prompt": "Final check",
        "votes": [{"persona_id": "p1", "vote": "support", "reason": "works"},
                  {"persona_id": "p2", "vote": "conditional", "reason": "channel missing"},
                  {"persona_id": "p3", "vote": "skeptical", "reason": "wallboard decides"}],
        "statements": [
            {"persona_id": "p1", "text": "It calculates live.", "stance": {"value": 2, "label": "support"}},
            {"persona_id": "p2", "text": "Slack is missing.", "stance": {"value": 1, "label": "conditional"}},
            {"persona_id": "p3", "text": "Not my lever.", "stance": {"value": -1, "label": "skeptical"}},
        ],
    }
    store.insert_council_session(council)
    return services.record_synthesis(
        "Final solution presentation", "hmw", ["c1"],
        {
            "gesamtbild": _BIG,
            "positionierung": ("For desk workers with calendar autonomy whose lunch is the only "
                               "meeting without an advocate: the window advocate makes the lunch "
                               "window as loud as its disruptors, slides visibly instead of dying "
                               "silently, computes counter-proposals from the real day and belongs "
                               "to the person alone — never to the reporting chain above, not a "
                               "wellness ideal, not a KPI, a negotiation tool for the one meeting "
                               "without a lobby, which is why the team rule beats every individual "
                               "justification and the care framing is never acceptable without comment."),
            "findings": [
                {"kind": "pain_solver",
                 "text": "DYNAMIC COUNTER-PROPOSAL: productises the only successful defence observed "
                         "in the field (a 13:45 counter-proposal accepted without comment)."},
                {"kind": "recommendation",
                 "text": "BUILD-SPEC CRITERION 1 — Calendar & channel binding: disruptions must come "
                         "from calendar, inbox and Slack (Fabian protosession_14a1aa34d5562f24, "
                         "verified live).",
                 "score": {"effort": 4, "value": 5}},
                {"kind": "recommendation",
                 "text": "CRITERION 2 — Honest failure: when the window hits a hard wall at 15:00 the "
                         "day is marked as lost, no auto-slide.",
                 "score": {"effort": 2, "value": 4}},
                {"kind": "open_question",
                 "text": "SURVEILLANCE TIPPING POINT (biggest risk): all five name the same boundary — "
                         "the moment window data becomes a KPI, acceptance flips to rejection."},
            ],
            "statements": [
                {"persona_id": "p1",
                 "text": "The counter-proposal is **provably** computed live from my current window — "
                         "my `lo-fi` stress test, passed this time.",
                 "stance": {"value": 2, "label": "support"},
                 "refs": [{"kind": "council", "id": "c1",
                           "quote": "My break rarely tips with a bang, it gets deprioritised."}],
                 "meta": {"persona_name": "Fabian Drees", "segment": "Calendar-autonomous desk workers"}},
                {"persona_id": "p3",
                 "text": "The wallboard decides, not me — an individual advocate changes nothing about "
                         "two outages in the early shift.",
                 "stance": {"value": -1, "label": "skeptical"},
                 "refs": [],
                 "meta": {"persona_name": "Janine Wolf", "segment": "Shift operations (non-target)"}},
            ],
        },
        store=store)


def _slide_model(syn, store):
    return _analytic_slides(syn, store, _SYNTHESIS_EXPORT_LABELS["en"], False,
                            "Final solution presentation", "Report")


def _frame_texts(prs):
    """[(slide_index, frame_text), …] over every text frame of the rendered deck."""
    out = []
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame:
                out.append((i, "\n".join(p.text for p in sh.text_frame.paragraphs)))
    return out


# ---------------------------------------------------------------- slide model


def test_deck_uses_master_layout_vocabulary(store):
    """Cover → verdict STATEMENT → takeaway cards → sentiment/stance CHARTS → quote + voices →
    finding cards → numbered recommendation cards + effort·value map: the report renders through
    the template vocabulary, not as content/prose slides."""
    syn = _showcase(store)
    kinds = [s["kind"] for s in _slide_model(syn, store)]
    assert kinds[0] == "cover"
    for kind in ("insight", "summary", "charts", "quote", "voices", "chart"):
        assert kind in kinds, f"missing {kind!r} slide (got {kinds})"


def test_verdict_statement_is_clamped_large_type(store):
    syn = _showcase(store)
    verdict = next(s for s in _slide_model(syn, store) if s["kind"] == "insight")
    assert verdict["eyebrow"] == "Verdict"
    assert 0 < len(verdict["statement"]) <= 240          # one sentence, large type — never a wall
    assert "THE WINNING CONCEPT:" not in verdict["statement"]   # the caps label is structure, not copy


def test_exec_summary_takeaways_from_authored_labels(store):
    """The caps-label spine of the authored exec prose becomes the takeaway cards (renderer-side
    splitting — headline + takeaways instead of a paragraph dump)."""
    syn = _showcase(store)
    summary = next(s for s in _slide_model(syn, store) if s["kind"] == "summary"
                   and s["heading"] == "Executive summary")
    titles = [it["title"] for it in summary["items"]]
    assert "WHY IT WINS" in titles and "SURVEY BACKING (5 answers)" in titles
    assert all(len(it["text"]) <= 260 for it in summary["items"])


def test_sentiment_and_stance_render_as_chart_slides(store):
    """Votes + contribution stances chart with scale-ordered, zero-free categories."""
    syn = _showcase(store)
    charts = next(s for s in _slide_model(syn, store) if s["kind"] == "charts")
    pie, bar = (it["chart"] for it in charts["items"])
    assert pie["type"] == "pie" and bar["type"] == "bar"
    assert pie["categories"] == ["Support", "Conditional", "Skeptical"]   # no zero buckets (V3)
    assert pie["values"] == [1, 1, 1] and bar["values"] == [1, 1, 1]


def test_votes_donut_is_a_valid_native_chart(store):
    """Round-3 H5: the sentiment slide's Votes card carries a NATIVE doughnut graphicFrame whose
    chart part is schema-valid — exactly ONE c:holeSize (the appended duplicate made strict
    renderers drop the chart: the 'legend with NO donut' empty panel) — with a nonzero series
    and every point carrying its semantic stance colour."""
    from pptx.oxml.ns import qn

    syn = _showcase(store)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx(syn["id"], store=store)))
    # the charts slide is identifiable by its 'Votes' card title; it must hold the graphicFrame
    votes_slide = next(sl for sl in prs.slides
                       if any(sh.has_text_frame and sh.text_frame.text == "Votes" for sh in sl.shapes))
    charts = [sh.chart for sh in votes_slide.shapes if getattr(sh, "has_chart", False)]
    assert len(charts) == 1, "the votes donut graphicFrame is missing from the sentiment slide"
    doughnuts = charts[0]._chartSpace.findall(".//" + qn("c:doughnutChart"))
    assert len(doughnuts) == 1
    holes = doughnuts[0].findall(qn("c:holeSize"))
    assert len(holes) == 1 and holes[0].get("val") == "62", (
        f"doughnut must carry exactly one c:holeSize (got {[h.get('val') for h in holes]}) — "
        "a duplicate makes the chart part invalid and PowerPoint renders an empty panel (H5)")
    series = charts[0].plots[0].series[0]
    assert sum(series.values) > 0 and list(series.values) == [1, 1, 1]
    assert len(doughnuts[0].findall(".//" + qn("c:dPt"))) == 3   # per-point stance colours applied


def test_voices_are_quote_slides_two_per_slide(store):
    syn = _showcase(store)
    slides = _slide_model(syn, store)
    quote = next(s for s in slides if s["kind"] == "quote")
    assert quote["attribution"] == "Fabian Drees"
    assert "deprioritised" in quote["text"]
    voices = [s for s in slides if s["kind"] == "voices"]
    assert voices and all(len(s["items"]) <= 2 for s in voices)
    fabian = voices[0]["items"][0]
    assert fabian["sentiment"] == "support" and fabian["sentiment_label"] == "Support"


def test_recommendations_are_numbered_cards_with_quiet_meta(store):
    syn = _showcase(store)
    rec_cards = next(s for s in _slide_model(syn, store)
                     if s["kind"] == "summary" and s.get("heading") == "Recommendations")
    assert rec_cards["items"][0]["title"].startswith("01 · ")
    assert rec_cards["items"][0]["meta"] == "Effort 4/5 · Value 5/5"
    # the effort·value map keeps SHORT legend labels (the heads), never the full prose
    rec_map = next(s for s in _slide_model(syn, store) if s["kind"] == "chart")
    assert all(len(p["label"]) <= 90 for p in rec_map["chart"]["points"])


# ------------------------------------------------------------- rendered deck


def test_no_rendered_frame_exceeds_the_prose_budget(store):
    """The §9 bar: no slide body is a text wall — every text frame stays within ~6 lines of
    continuous prose (≤700 chars per frame, ≤500 per paragraph)."""
    syn = _showcase(store)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx(syn["id"], store=store)))
    for i, text in _frame_texts(prs):
        assert len(text) <= 700, f"slide {i + 1}: frame carries {len(text)} chars of prose:\n{text[:200]}"
        for para in text.split("\n"):
            assert len(para) <= 500, f"slide {i + 1}: paragraph runs {len(para)} chars"


def test_no_markdown_artifacts_or_bare_ids_on_slides(store):
    syn = _showcase(store)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx(syn["id"], store=store)))
    deck_text = " ".join(t for _, t in _frame_texts(prs))
    assert "**" not in deck_text and "`" not in deck_text
    assert "protosession_" not in deck_text          # terminal ids are not slide copy
    assert "provably" in deck_text                   # the de-markdowned content survived


def test_project_section_prose_is_budgeted_with_report_footnote(store):
    """Project-report sections clamp to the slide budget and point to the full report instead of
    dumping the section body."""
    long_para = "The full analysis goes on and on with detail after detail. " * 40
    rep = {"id": "rlong", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-10T00:00:00+00:00", "lead": "", "council_ids": [],
           "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": "s1", "heading": "Findings", "markdown": long_para,
                         "citations": [], "source_study_ids": [], "figures": []}]}
    store.upsert_synthesis(rep)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx("rlong", store=store)))
    texts = [t for _, t in _frame_texts(prs)]
    assert all(len(t) <= 700 for t in texts)
    assert any("Details in the full report" in t for t in texts)


def test_project_deck_translates_engine_phases_into_stakeholder_chapters(store):
    phases = [
        ("Product understanding", "The comparison covers two customer-facing variants."),
        ("Cohort integrity", "Eight simulated participants supplied directional reactions."),
        ("React", "Variant B was preferred by six of eight participants."),
        ("Gate", "Decision: continue with variant B after making the interruption risk visible."),
    ]
    rep = {
        "id": "rpresent", "title": "Concept comparison — Report", "scope": "project",
        "project_id": "", "created_at": "2026-08-22T00:00:00+00:00",
        "lead": ("This report carries the evidence along the research phases from the initial "
                 "question to prioritised conclusions."),
        "council_ids": [], "findings": [], "statements": [], "prompts": [],
        "graph_snapshot": None,
        "sections": [{
            "id": f"s{index}", "heading": heading, "markdown": markdown,
            "intent": f"Author the {heading} phase ({'converge' if index == 4 else 'diverge'}) "
                      "grounded in its evidence + what it produced.",
            "citations": [], "source_study_ids": ["note:a", "council:b"], "figures": [],
        } for index, (heading, markdown) in enumerate(phases, 1)],
    }
    store.upsert_synthesis(rep)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx("rpresent", store=store)))
    deck_text = "\n".join(text for _, text in _frame_texts(prs))

    for label in ("Test subject & context", "Participants & limitations",
                  "Reactions & key findings", "Decision & recommendations"):
        assert label in deck_text
    lines = {line.strip() for line in deck_text.splitlines()}
    assert not any(heading in lines for heading, _markdown in phases)
    assert "Decision: continue with variant B" in deck_text  # result-led cover subtitle
    assert "2 sources in the full report" in deck_text
    assert "note:a" not in deck_text and "council:b" not in deck_text


# ------------------------------------------------------------------- helpers


def test_split_card_never_breaks_words_or_times():
    head, body = _split_card("Kalender-autonome Schreibtisch-Arbeiter:innen (Zielsegment) — Stance: win")
    assert head == "Kalender-autonome Schreibtisch-Arbeiter:innen (Zielsegment)"
    assert body == "Stance: win"
    head, body = _split_card("Hard walls at 15:00 stay fixed: the window yields.")
    assert head == "Hard walls at 15:00 stay fixed" and body == "the window yields."
    assert _split_card("No separator here at all") == ("", "No separator here at all")


def test_clamp_prose_cuts_at_sentence_and_clause_boundaries():
    text = "First point stands. Second point follows. Third point would overflow the budget."
    out, truncated = _clamp_prose(text, 45)
    assert out == "First point stands. Second point follows." and truncated
    out, truncated = _clamp_prose("One enormous clause, with a second clause, and a third clause", 40)
    assert out.endswith("…") and len(out) <= 45 and truncated


def test_label_segments_finds_the_caps_spine():
    segs = _label_segments(_BIG)
    assert [s[0] for s in segs] == ["THE WINNING CONCEPT", "WHY IT WINS", "SURVEY BACKING (5 answers)"]
