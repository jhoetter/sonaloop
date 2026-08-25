"""Evidence-linked presentation plans: gather → author → persist."""
from __future__ import annotations

import io

import pytest
from pptx import Presentation

from sonaloop import services
from sonaloop.models import Synthesis
from sonaloop.presentation import (
    PRESENTATION_PLAN_SCHEMA,
    presentation_plan_qa,
    validate_presentation_plan,
)
from sonaloop.web._report import render_report
from conftest import create_persona


PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


def _report(store):
    project = services.start_project(
        "Reaction deck", "Which stimulus should ship?", methodology="Reaction Test",
        operation_id="presentation-test-project", store=store,
    )
    report = Synthesis(
        id="report_presentation_test", title="Reaction deck — Report", start_input="",
        council_ids=[], arc_narrative="", gesamtbild="", positionierung="", references=[],
        created_at="2026-08-22T00:00:00+00:00",
        scope="project", project_id=project["id"], status="done",
        lead="Use B after one important revision.",
        sections=[{
            "id": "decision", "heading": "Decision", "markdown": "Use B, then revise.",
            "source_study_ids": [], "citations": [], "figures": [], "status": "done",
        }],
    ).to_dict()
    store.upsert_synthesis(report)
    return project, report


def _plan():
    notes = {"talk_track": "Explain the evidence behind this slide.",
             "caveats": ["Synthetic cohort, not observed customer behavior."]}
    return {
        "schema": PRESENTATION_PLAN_SCHEMA,
        "title": "Reaction decision", "audience": "stakeholder",
        "objective": "Choose the next action", "duration_minutes": 8,
        "slides": [
            {"id": "cover", "kind": "cover", "headline": "B is the stronger base",
             "speaker_notes": notes},
            {"id": "result", "kind": "stats", "headline": "Preference moves from 6/8 to 8/8",
             "items": [{"label": "Blind", "value": "6/8"},
                       {"label": "With context", "value": "8/8"}],
             "evidence_refs": ["synthesis:decision"], "speaker_notes": notes},
        ],
        "appendix": [
            {"id": "sources", "kind": "source_index", "headline": "Sources",
             "speaker_notes": {"talk_track": "Use this for questions."}},
        ],
    }


def test_methodology_exposes_data_authored_deck_profile(store):
    profile = services.get_methodology("reaction_test", store=store)["presentation"]["deck"]
    assert profile["purpose"]
    assert {row["role"] for row in profile["story_beats"]} >= {
        "decision", "stimulus", "cohort", "reaction", "revision",
    }
    assert "6–8" in profile["target_core_slides"]
    assert "response table by persona" in profile["appendix"]


def test_presentation_plan_validation_and_qa():
    plan = validate_presentation_plan(_plan())
    assert plan["slides"][1]["speaker_notes"]["takeaway"] == \
        "Preference moves from 6/8 to 8/8"
    qa = presentation_plan_qa(plan)
    assert qa["speaker_notes_count"] == 3
    assert qa["core_slide_count"] == 2
    broken = _plan()
    broken["slides"][1]["evidence_refs"] = []
    with pytest.raises(ValueError, match="evidence_refs"):
        validate_presentation_plan(broken)

    verbose = _plan()
    verbose["slides"][0]["headline"] = (
        "This unusually long presentation headline cannot be understood or fitted at a glance"
    )
    warnings = presentation_plan_qa(validate_presentation_plan(verbose))["warnings"]
    assert any(row["code"] == "headline_too_long" and row["slide_id"] == "cover"
               for row in warnings)


def test_brief_and_retry_safe_record(store):
    project, report = _report(store)
    brief = services.brief_presentation(report["id"], duration_minutes=12, store=store)
    assert brief["project_id"] == project["id"]
    assert brief["methodology"]["deck_profile"]["story_beats"]
    assert brief["methodology"]["result_contract"]["schemas"][0]["id"] == \
        "stimulus_reaction.v1"
    assert brief["output_contract"]["schema"] == PRESENTATION_PLAN_SCHEMA
    assert "decision_dashboard" in brief["output_contract"]["preferred_blueprints"]

    first = services.record_presentation_plan(
        report["id"], _plan(), operation_id="deck-v1", store=store)
    replay = services.record_presentation_plan(
        report["id"], _plan(), operation_id="deck-v1", store=store)
    assert first["presentation_plan_revision"] == 1
    assert replay["idempotent_replay"] is True
    changed = _plan()
    changed["title"] = "Different"
    with pytest.raises(ValueError, match="PRESENTATION_OPERATION_CONFLICT"):
        services.record_presentation_plan(
            report["id"], changed, operation_id="deck-v1", store=store)


def test_every_packaged_methodology_has_a_deck_prompt_profile(store):
    for row in services.list_methodologies(store=store):
        profile = services.get_methodology(row["key"], store=store).get("presentation", {}).get("deck")
        assert profile and profile["story_beats"], row["key"]
        assert profile["required_visuals"] and profile["appendix"] and profile["avoid"]


def test_stored_plan_renders_visual_story_native_notes_and_appendix(store):
    project, report = _report(store)
    persona_ids = [create_persona(store, name) for name in ("Alba Costa", "Bruno Keller")]
    screen = services.attach_asset(
        project["id"], content_base64=PNG_BASE64, filename="stimulus-b.png",
        title="Variant B", store=store,
    )
    notes = {
        "takeaway": "B wins only after one revision.",
        "talk_track": "Explain what changed between the blind and informed reactions.",
        "evidence": ["council:round-2"],
        "caveats": ["Synthetic reactions are directional evidence."],
        "transition": "Now show the exact revision.",
        "timing_seconds": 45,
    }
    plan = {
        "schema": PRESENTATION_PLAN_SCHEMA,
        "title": "Reaction decision", "audience": "project team",
        "objective": "Choose the next iteration", "duration_minutes": 10,
        "slides": [
            {"id": "cover", "kind": "cover", "headline": "B is the stronger base",
             "subheadline": "One critical revision remains", "speaker_notes": notes},
            {"id": "stimuli", "kind": "stimulus_comparison", "headline": "What participants saw",
             "left": {"label": "A", "asset_id": screen["id"], "callouts": ["More detail"]},
             "right": {"label": "B", "asset_id": screen["id"], "highlight": True,
                       "callouts": ["Clearer action"]},
             "evidence_refs": [screen["id"]], "speaker_notes": notes},
            {"id": "cohort", "kind": "persona_grid", "headline": "Two distinct customer lenses",
             "items": [{"persona_id": persona_ids[0], "badge": "digital"},
                       {"persona_id": persona_ids[1], "badge": "cautious"}],
             "evidence_refs": persona_ids, "speaker_notes": notes},
            {"id": "movement", "kind": "preference_shift", "headline": "Context resolves the split",
             "before": {"label": "Blind reaction", "value": 6, "total": 8,
                        "detail": "B preferred"},
             "after": {"label": "After explanation", "value": 8, "total": 8,
                        "detail": "B preferred"},
             "switchers": [{"persona_id": persona_ids[1], "reason": "Context resolved the risk"}],
             "evidence_refs": ["council:round-2"],
             "speaker_notes": notes},
            {"id": "revision", "kind": "revision_mockup", "headline": "Make the interruption risk explicit",
             "asset_id": screen["id"],
             "proposal": {"eyebrow": "Before you start", "headline": "Plan ten uninterrupted minutes",
                          "body": "If you leave the process, the switch restarts.",
                          "primary_cta": "Start switch", "secondary_cta": "Later"},
             "why": ["Sets expectation", "Names the restart consequence"],
             "evidence_refs": ["council:round-2"], "speaker_notes": notes},
            {"id": "risk", "kind": "decision_dashboard",
             "headline": "B should proceed only with one revision",
             "decision": {"label": "Revise", "text": "Use B as the base",
                          "detail": "Do not ship the interruption copy as-is."},
             "metrics": [{"value": "8/8", "label": "Prefer B after context"},
                         {"value": "1", "label": "Critical copy gap"}],
             "rationale": [{"title": "Keep", "text": "Clear preparation and timing."},
                           {"title": "Change", "text": "Name interruption and restart risk."},
                           {"title": "Validate", "text": "Observe completion behavior."}],
             "evidence_refs": ["report:limitations"], "speaker_notes": notes},
            {"id": "next", "kind": "next_steps", "headline": "Revise, test, decide",
             "steps": [{"label": "Now", "title": "Revise copy", "text": "Add duration and restart risk."},
                       {"label": "Next", "title": "Behavioral test", "text": "Observe completion and confusion."}],
             "evidence_refs": ["report:decision"], "speaker_notes": notes},
        ],
        "appendix": [
            {"id": "persona-detail", "kind": "persona_detail", "headline": "Persona detail",
             "items": [{"persona_id": persona_ids[0], "quote": "I need a clear start."},
                       {"persona_id": persona_ids[1], "quote": "Tell me what happens next."}],
             "speaker_notes": notes},
            {"id": "sources", "kind": "source_index", "headline": "Sources",
             "columns": ["Source", "Contribution"],
             "rows": [["Reaction round 2", "Preference after context"]],
             "speaker_notes": notes},
        ],
    }
    services.record_presentation_plan(
        report["id"], plan, operation_id="visual-deck-v1", store=store)
    assert presentation_plan_qa(validate_presentation_plan(plan))["status"] == "pass"

    data = services.export_synthesis_pptx(report["id"], store=store)
    deck = Presentation(io.BytesIO(data))
    visible = "\n".join(
        shape.text_frame.text for slide in deck.slides for shape in slide.shapes
        if shape.has_text_frame
    )

    assert len(deck.slides) == 10  # seven core + appendix divider + two appendix slides
    assert "What participants saw" in visible
    assert "Alba Costa" in visible and "Bruno Keller" in visible
    assert "I need a clear start." in visible
    assert "6/8" in visible and "8/8" in visible
    assert "Use B as the base" in visible and "Plan ten uninterrupted minutes" in visible
    assert "Start switch" in visible and "Bruno Keller" in visible
    assert "Evidence, method and sources" in visible
    assert "Thank you" not in visible
    assert all("TALK TRACK" in slide.notes_slide.notes_text_frame.text
               for slide in deck.slides)
    assert "45 SEC" in deck.slides[0].notes_slide.notes_text_frame.text
    assert "Synthetic reactions are directional evidence." in \
        deck.slides[0].notes_slide.notes_text_frame.text

    html = str(render_report(store.get_synthesis(report["id"]), store, audience="stakeholder"))
    assert "delivery-report" in html
    assert "What participants saw" in html
    assert html.count(screen["url"]) >= 2
    assert "Alba Costa" in html and "Bruno Keller" in html
    assert "Plan ten uninterrupted minutes" in html
    assert html.rfind("Method and interpretation") > html.find("Sources")
